"""Talentrupt presentation builder (.pptx) via python-pptx.

Each deck is laid out by an LLM design plan (a layout chosen PER slide from a varied
set) and rendered by distinct layout renderers — so slides differ and read as designed,
not a single repeated template. The cover can use an AI-generated, brand-grounded
background image. Text + design are grounded in the ingested brand library.
"""
from __future__ import annotations

import random
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from ..knowledge import retrieve
from ..models import Brand, Campaign
from ..providers import llm
from .common import logo_path, public_url, storage_subdir, unique_name

NAVY = RGBColor(0x0B, 0x35, 0x59)
RED = RGBColor(0xF6, 0x40, 0x4C)
CORAL = RGBColor(0xFF, 0x7A, 0x52)
CREAM = RGBColor(0xEB, 0xE9, 0xDF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SW, SH = Inches(13.333), Inches(7.5)
HEAD = "Segoe UI Semibold"
BODY = "Segoe UI"

LAYOUTS = {"cover", "section", "bullets", "metric", "two_column", "quote", "comparison", "closing",
           "cards", "timeline"}

# design_theme → accent profile. Navy/red/cream + logo are ALWAYS on-brand; the theme only re-weights
# the accent colour, the side-spine width, and the chip colour so decks vary without going off-brand.
DECK_THEMES = {
    "editorial": {"accent": RED, "spine": 0.30, "chip": RED},      # default — matches today's look
    "minimal": {"accent": NAVY, "spine": 0.16, "chip": NAVY},
    "bold": {"accent": RED, "spine": 0.46, "chip": CORAL},
    "data-driven": {"accent": RED, "spine": 0.30, "chip": CORAL},
}


def _deck_theme(name: str) -> dict:
    return DECK_THEMES.get((name or "").strip().lower(), DECK_THEMES["editorial"])


# Voice per tone, woven into the outline prompt so the narrative matches what the user asked for.
_TONE_VOICE = {
    "executive": "boardroom-level: concise, confident, lead every slide with the outcome",
    "professional": "polished, credible US B2B (the default)",
    "conversational": "direct and warm — speak to the reader as 'you'",
    "persuasive": "sales-forward — build toward a strong close and a clear ask",
    "data-driven": "tight and evidence-led — foreground the single most relevant REAL proof point",
}
_DEPTH_VOICE = {
    "concise": "headline-driven; at most 3 short bullets per slide; lots of whitespace",
    "standard": "balanced detail",
    "in-depth": "thorough supporting detail and a process/cards slide where it helps",
}


def _directives(audience: str, tone: str, depth: str) -> str:
    """Build the per-brief directive block injected into the outline prompt (only for set values)."""
    lines = []
    if audience:
        lines.append(f"Write for this audience: {audience}. Frame the problem, examples and CTA for them.")
    if tone and tone in _TONE_VOICE:
        lines.append(f"Voice/tone: {_TONE_VOICE[tone]}.")
    if depth and depth in _DEPTH_VOICE:
        lines.append(f"Depth: {_DEPTH_VOICE[depth]}.")
    return ("\nBRIEF DIRECTIVES:\n- " + "\n- ".join(lines) + "\n") if lines else ""


# --- Outline / design plan ------------------------------------------------
async def _outline(brand: Brand | None, campaign: Campaign | None, topic: str, n: int, *,
                   audience: str = "", tone: str = "", depth: str = "", design_theme: str = "editorial") -> list[dict]:
    if llm.provider_available():
        pillars = ", ".join(brand.pillars) if brand and brand.pillars else ""
        proof = "; ".join(brand.proof_points) if brand and brand.proof_points else ""
        services = ", ".join(brand.services) if brand and brand.services else ""
        context = await retrieve.brand_context(topic, k=6)
        sys = (
            "You are Talentrupt's strategy deck designer (offshore RPO selling into the US market, "
            f"'RPO Done Right'). Brand pillars: {pillars}. Services: {services}. "
            f"Proof points — these are the ONLY real numbers you may ever state: {proof}.\n"
            + (f"\n{context}\n\n" if context else "")
            + _directives(audience, tone, depth)
            + f"Design a {n}-slide deck built SPECIFICALLY around the topic below — its own narrative "
            "arc, not a generic company overview. Lead with the audience's problem or the topic's core "
            "idea, then show how Talentrupt's services solve it. Reference the audience/client and topic "
            "by name where natural, and weave in real Talentrupt services and proof so it feels tailored, "
            "not boilerplate.\n"
            "EVERY slide must be SUBSTANTIVE — no near-empty slides. Fill each with real, specific content "
            "(not one line floating in space).\n"
            "ANTI-FABRICATION (critical): do NOT state ANY percentage, multiple, count or dollar figure "
            "ANYWHERE — titles, subtitles, bodies, bullets, card details, step descriptions, notes — "
            f"UNLESS it is copied EXACTLY from this list: {proof or 'none'}. Otherwise describe benefits "
            "QUALITATIVELY (write 'faster time-to-fill', NEVER '40% faster'; 'lower overhead', not '30% "
            "cheaper'). Inventing a number is a serious error.\n"
            f"Return ONLY JSON: {{\"slides\": [...]}} with {n} items. Each slide object:\n"
            '- "layout": one of "cover","section","bullets","metric","two_column","quote",'
            '"comparison","cards","timeline","closing"\n'
            '- "title": short, executive, specific to the topic\n'
            '- "subtitle": one supporting line — REQUIRED for "cover","section","closing","timeline"\n'
            '- "bullets": 3-4 SUBSTANTIVE points — each a specific, benefit-driven sentence with a '
            'concrete detail tied to the audience, not a generic phrase (for "bullets")\n'
            '- "metric": JUST a short number/percent copied from the proof points — e.g. "90%", "500+" — '
            '6 chars MAX, NEVER a sentence; put words in "metric_label" (for "metric"). AT MOST one metric '
            "slide, only with a proof point relevant to the topic\n"
            '- "quote": a punchy one-sentence pull-quote + "attribution" (for "quote")\n'
            '- "left"/"right": {"label","points":[...3-4 points...]} (for "two_column"/"comparison")\n'
            '- "cards": [{"title","detail"}] exactly 3 value props, each with a FULL one-line detail '
            '(for "cards")\n'
            '- "steps": [{"label": short, "desc": a 6-12 word benefit line}] 3-5 steps (for "timeline") — '
            "every step MUST have a real 'desc', never a bare label\n"
            '- "notes": 2-3 sentences of speaker talking points for this slide\n'
            "LAYOUT MIX (so the deck looks designed, never an all-bullets wall): slide 1 MUST be 'cover' "
            "and the LAST MUST be 'closing'; include AT LEAST one 'section' divider and at least one of "
            "'two_column'/'comparison'/'cards'/'timeline'; use AT MOST one 'metric' and one 'quote'; NO "
            "MORE THAN half the slides may be 'bullets'. Professional US B2B tone, no filler.\n"
            "Example arc: cover → section → bullets(the problem, for THIS client) → cards(how Talentrupt "
            "helps) → metric(one real proof) → timeline(Source/Screen/Submit/Place, each with a desc) → "
            "closing(CTA)."
        )
        usr = f"Topic: {topic}" + (
            f" | Campaign: {campaign.name}, audience: {campaign.audience}" if campaign else ""
        )
        # Slightly higher temperature so repeated requests on a topic diverge rather than repeat.
        try:
            data = await llm.chat_json(
                [{"role": "system", "content": sys}, {"role": "user", "content": usr}], temperature=0.85
            )
            slides = data.get("slides") if isinstance(data, dict) else None
            if slides:
                return _normalize(slides[:n], topic)
        except Exception:
            pass  # transient provider/parse error -> degrade to the varied fallback below

    return _fallback_outline(brand, campaign, topic, n)


def _fallback_outline(brand: Brand | None, campaign: Campaign | None, topic: str, n: int) -> list[dict]:
    """Deterministic but VARIED fallback (no provider): pick one slide per role from small pools so
    repeat builds differ instead of emitting one identical template. Any number is a REAL proof point."""
    rnd = random.Random()  # system entropy -> distinct selection each build
    title = topic or (campaign.name if campaign else "Talentrupt RPO")
    proofs = (brand.proof_points if brand and brand.proof_points
              else ["90% submission-to-interview alignment", "500+ healthcare roles filled annually"])
    challenge_pool = [
        ["Recruiting demand outpaces internal capacity", "High overhead of in-house recruiting",
         "Inconsistent submission quality and SLA misses"],
        ["Time-to-fill keeps slipping on critical roles", "Hiring spikes overwhelm a fixed team",
         "Thin pipeline visibility and reporting"],
        ["Quality candidates surface too slowly", "Cost-per-hire creeps up every quarter",
         "Coverage gaps across roles and time zones"],
    ]
    cards_pool = [
        [{"title": "Scalable", "detail": "Capacity that flexes with demand"},
         {"title": "SLA-driven", "detail": "On-time, predictable delivery"},
         {"title": "Transparent", "detail": "Live reporting end to end"}],
        [{"title": "Dedicated", "detail": "Recruiters embedded in your team"},
         {"title": "Cost-efficient", "detail": "Lower overhead than in-house"},
         {"title": "Fast ramp", "detail": "Productive in days, not months"}],
    ]
    closing_pool = [
        {"title": "Let's scale your hiring", "subtitle": "Partner with Talentrupt — RPO Done Right."},
        {"title": "Ready when you are", "subtitle": "Let's build your hiring engine — RPO Done Right."},
    ]
    out = [
        {"layout": "cover", "title": title, "subtitle": "RPO Done Right — engineered hiring systems"},
        {"layout": "section", "title": "The Challenge", "subtitle": "Where hiring breaks down"},
        {"layout": "bullets", "title": "What's slowing hiring down", "bullets": rnd.choice(challenge_pool)},
        {"layout": "cards", "title": "Why Talentrupt works", "cards": rnd.choice(cards_pool)},
        {"layout": "metric", "title": "Proven delivery", "metric": rnd.choice(proofs)},
        {"layout": "timeline", "title": "How we deliver", "subtitle": "A proven, SLA-driven process end to end",
         "steps": [{"label": "Source", "desc": "Pinpoint the right talent, fast"},
                   {"label": "Screen", "desc": "Vet for skills, fit and quality"},
                   {"label": "Submit", "desc": "Deliver shortlists on SLA"},
                   {"label": "Place", "desc": "Close, onboard and support"}]},
        {"layout": "comparison", "title": "RPO vs In-House",
         "left": {"label": "Talentrupt RPO", "points": ["Scales on demand", "SLA-driven", "Lower overhead"]},
         "right": {"label": "In-house only", "points": ["Capacity-bound", "Slower ramp", "Higher cost"]}},
    ]
    rnd.shuffle(closing_pool)
    closing = {**closing_pool[0], "layout": "closing"}
    body = out[1:]  # keep cover first
    rnd.shuffle(body)  # vary the middle order run-to-run
    out = [out[0]] + body
    out = out[:max(2, n - 1)] + [closing]
    return out[:n]


def _normalize(slides: list, topic: str) -> list[dict]:
    out = []
    for i, s in enumerate(slides):
        if not isinstance(s, dict):  # tolerate a stray non-object slide from the model
            s = {}
        layout = s.get("layout") if s.get("layout") in LAYOUTS else "bullets"
        if i == 0:
            layout = "cover"
        out.append({**s, "layout": layout})
    if out and out[-1]["layout"] not in ("closing", "cover"):
        out[-1]["layout"] = "closing"
    return out


# --- Low-level helpers ----------------------------------------------------
def _bg(slide, color):
    _rect(slide, 0, 0, SW, SH, color)


def _rect(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _round(slide, x, y, w, h, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color; shp.line.fill.background()
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = 0.08
    except Exception:
        pass
    return shp


def _text(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT, anchor=None, font=None):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font or (HEAD if bold else BODY)
    return box


def _bullet_list(slide, x, y, w, points, color, size=18, gap=0.52):
    for j, pt in enumerate(points[:5]):
        yy = y + Inches(j * gap)
        _rect(slide, x, yy + Inches(0.1), Inches(0.12), Inches(0.12), RED)
        _text(slide, x + Inches(0.32), yy, w - Inches(0.32), Inches(0.5), pt, size, color)


def _logo(slide, color):
    # Real Talentrupt logo (navy TR in a red square) + wordmark; text-only fallback if it can't load.
    try:
        slide.shapes.add_picture(str(logo_path()), Inches(0.55), Inches(0.28), height=Inches(0.6))
        _text(slide, Inches(1.3), Inches(0.4), Inches(3), Inches(0.4), "TALENTRUPT", 14, color, bold=True)
    except Exception:
        _text(slide, Inches(0.6), Inches(0.34), Inches(3), Inches(0.4), "TR  TALENTRUPT", 14, color, bold=True)


def _footer(slide, idx, total, color=NAVY):
    _text(slide, Inches(0.6), Inches(7.06), Inches(6), Inches(0.3), "Talentrupt · RPO Done Right", 10, color)
    _text(slide, Inches(11.6), Inches(7.06), Inches(1.1), Inches(0.3), f"{idx}/{total}", 10, color, align=PP_ALIGN.RIGHT)


# --- AI cover background --------------------------------------------------
async def _ai_cover(brand: Brand | None, topic: str) -> str | None:
    if not llm.image_provider_available():
        return None
    try:
        from . import images
        refs = images._load_references(await retrieve.image_references(topic, n=2))
        prompt = (
            "A premium, abstract presentation COVER background for Talentrupt, an offshore RPO "
            "(recruitment) firm. Brand colors: deep navy #0B3559, coral red #F6404C accents, warm "
            f"cream #EBE9DF. Modern editorial/magazine feel evoking: {topic}. Subtle geometric shapes, "
            "soft depth, professional. Leave the lower-left area relatively clear. NO text, NO words, "
            "NO logos. Wide 3:2."
            + ("\nMatch the brand palette and finish of the attached Talentrupt references, original composition." if refs else "")
        )
        if refs:
            data = await llm.generate_image_edit(prompt, refs, size="1536x1024")
        else:
            data = await llm.generate_image_bytes(prompt, size="1536x1024")
        fname = unique_name("tr-cover", "png")
        path = storage_subdir("images") / fname
        with open(path, "wb") as f:
            f.write(data)
        return str(path)
    except Exception:
        return None


# --- Slide renderers ------------------------------------------------------
def _r_cover(slide, s, cover_img, th):
    title = s.get("title", "Talentrupt")
    sub = s.get("subtitle") or "RPO Done Right"
    if cover_img:
        try:
            slide.shapes.add_picture(cover_img, 0, 0, width=SW, height=SH)
        except Exception:
            _bg(slide, CREAM)
        _rect(slide, 0, Inches(5.0), SW, Inches(2.5), NAVY)  # legibility band
        _rect(slide, 0, Inches(4.92), SW, Inches(0.08), th["accent"])
        _logo(slide, CREAM)
        _text(slide, Inches(0.7), Inches(5.15), Inches(12), Inches(1.4), title, 44, WHITE, bold=True)
        _text(slide, Inches(0.72), Inches(6.55), Inches(12), Inches(0.6), sub, 18, CREAM)
    else:
        _bg(slide, CREAM)
        _rect(slide, 0, 0, Inches(th["spine"]), SH, th["accent"])
        _logo(slide, NAVY)
        _text(slide, Inches(0.9), Inches(2.6), Inches(11.4), Inches(2.4), title, 48, NAVY, bold=True)
        _rect(slide, Inches(0.95), Inches(4.95), Inches(2.2), Inches(0.12), th["accent"])
        _text(slide, Inches(0.9), Inches(5.2), Inches(11), Inches(1), sub, 22, NAVY)


def _r_section(slide, s, th):
    _bg(slide, NAVY)
    _rect(slide, 0, 0, Inches(th["spine"]), SH, th["accent"])
    _logo(slide, CREAM)
    _text(slide, Inches(0.9), Inches(2.7), Inches(11.4), Inches(2), s.get("title", ""), 40, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    _rect(slide, Inches(0.95), Inches(4.6), Inches(2), Inches(0.12), th["accent"])
    if s.get("subtitle"):
        _text(slide, Inches(0.9), Inches(4.85), Inches(11), Inches(1), s["subtitle"], 20, CREAM)


def _r_bullets(slide, s, idx, total, th):
    _bg(slide, WHITE)
    _rect(slide, 0, 0, Inches(th["spine"]), SH, th["accent"])
    _rect(slide, 0, 0, SW, Inches(1.5), NAVY)
    _text(slide, Inches(0.7), Inches(0.55), Inches(11.5), Inches(0.8), s.get("title", ""), 30, WHITE, bold=True)
    shown = (s.get("bullets") or [])[:5]
    top, ch = 2.0, min(0.95, (6.7 - 2.0) / max(1, len(shown)))
    for n, b in enumerate(shown, 1):
        y = Inches(top + (n - 1) * ch)
        _round(slide, Inches(0.7), y, Inches(11.9), Inches(ch - 0.18), CREAM)
        _round(slide, Inches(0.95), y + Inches(0.14), Inches(0.5), Inches(0.5), th["chip"])
        _text(slide, Inches(0.95), y + Inches(0.18), Inches(0.5), Inches(0.4), str(n), 18, WHITE, bold=True, align=PP_ALIGN.CENTER)
        _text(slide, Inches(1.7), y + Inches(0.04), Inches(10.5), Inches(ch - 0.2), str(b), 17, NAVY, anchor=MSO_ANCHOR.MIDDLE)
    _footer(slide, idx, total)


def _r_cards(slide, s, idx, total, th):
    _bg(slide, WHITE)
    _rect(slide, 0, 0, Inches(th["spine"]), SH, th["accent"])
    _rect(slide, 0, 0, SW, Inches(1.5), NAVY)
    _text(slide, Inches(0.7), Inches(0.55), Inches(11.5), Inches(0.8), s.get("title", ""), 30, WHITE, bold=True)
    cards = [c for c in (s.get("cards") or []) if isinstance(c, dict)][:3]
    if not cards:
        return _footer(slide, idx, total)
    n = len(cards)
    gap = 0.4
    cw = (11.9 - gap * (n - 1)) / n
    for j, c in enumerate(cards):
        x = Inches(0.7 + j * (cw + gap))
        _round(slide, x, Inches(2.1), Inches(cw), Inches(4.3), CREAM)
        _round(slide, x + Inches(0.3), Inches(2.45), Inches(0.7), Inches(0.7), th["chip"])
        _text(slide, x + Inches(0.3), Inches(2.55), Inches(0.7), Inches(0.5), str(j + 1), 24, WHITE, bold=True, align=PP_ALIGN.CENTER)
        _text(slide, x + Inches(0.3), Inches(3.45), Inches(cw - 0.6), Inches(0.9), c.get("title", ""), 20, NAVY, bold=True)
        _text(slide, x + Inches(0.3), Inches(4.4), Inches(cw - 0.6), Inches(1.8), c.get("detail", ""), 15, NAVY)
    _footer(slide, idx, total)


def _r_timeline(slide, s, idx, total, th):
    _bg(slide, WHITE)
    _rect(slide, 0, 0, Inches(th["spine"]), SH, th["accent"])
    _rect(slide, 0, 0, SW, Inches(1.5), NAVY)
    _text(slide, Inches(0.7), Inches(0.42), Inches(11.5), Inches(0.7), s.get("title", ""), 30, WHITE, bold=True)
    sub = s.get("subtitle") or s.get("intro") or ""
    if sub:
        _text(slide, Inches(0.72), Inches(1.04), Inches(11.5), Inches(0.4), sub, 14, CREAM)
    # Steps may be {"label","desc"} or bare strings; normalize and always carry a description.
    norm = []
    for st_ in (s.get("steps") or [])[:5]:
        if isinstance(st_, dict):
            norm.append((str(st_.get("label", "") or "").strip(),
                         str(st_.get("desc") or st_.get("detail") or "").strip()))
        elif str(st_).strip():
            norm.append((str(st_).strip(), ""))
    if not norm:
        norm = [("Source", "Pinpoint the right talent, fast"), ("Screen", "Vet for fit and quality"),
                ("Submit", "Deliver shortlists on SLA"), ("Place", "Close, onboard and support")]
    n = len(norm)
    cw = 11.9 / n
    rowy = 2.75
    _rect(slide, Inches(1.0), Inches(rowy + 0.47), Inches(11.3), Inches(0.05), CREAM)  # connector behind chips
    for j, (label, desc) in enumerate(norm):
        cx = 0.7 + j * cw + cw / 2
        _round(slide, Inches(cx - 0.5), Inches(rowy), Inches(1.0), Inches(1.0), th["chip"])
        _text(slide, Inches(cx - 0.5), Inches(rowy + 0.2), Inches(1.0), Inches(0.6), str(j + 1), 32, WHITE, bold=True, align=PP_ALIGN.CENTER)
        _text(slide, Inches(0.7 + j * cw), Inches(rowy + 1.2), Inches(cw - 0.15), Inches(0.5), label, 19, NAVY, bold=True, align=PP_ALIGN.CENTER)
        if desc:
            _text(slide, Inches(0.7 + j * cw + 0.1), Inches(rowy + 1.78), Inches(cw - 0.4), Inches(1.7), desc, 12.5, NAVY, align=PP_ALIGN.CENTER)
    _footer(slide, idx, total)


def _metric_parts(s: dict) -> tuple[str, str]:
    """Split a metric into a SHORT headline number + a descriptive label, so the giant display
    text can never receive a long phrase (which overflows the slide). Models sometimes return the
    whole proof point as the metric — e.g. "500+ healthcare roles filled annually" becomes
    ("500+", "healthcare roles filled annually")."""
    raw = " ".join(str(s.get("metric", "") or "").split())
    label = str(s.get("metric_label", "") or "").strip()
    m = re.match(r"[~<>$€£]?\d[\d,.]*\s*(?:%|\+|x|×|k|m|b|bn|hrs?|days?|hours?)*\+?", raw, re.I)
    if m and m.group().strip():
        num = m.group().strip()
        rest = raw[m.end():].strip(" .,:;–—-")
        if rest and not label:
            label = rest  # carry the description into the label rather than dropping it
        # Keep the number INTACT (no truncation — that would show a wrong stat); _metric_size shrinks it.
        return num[:20], label
    # No leading number — keep the display short and push any overflow to the label.
    if len(raw) > 16 and not label:
        label = raw
    return raw[:16], label


def _metric_is_real(s: dict, brand: Brand | None) -> bool:
    """Anti-fabrication gate for the standout-number slide: the metric is shown ONLY if its number
    matches a REAL brand proof point (exact leading-number match), so the LLM can't slip a made-up
    figure (e.g. '40%') onto the one slide whose entire job is to display a statistic."""
    num, _ = _metric_parts(s)
    if not num:
        return False
    target = num.replace(" ", "").lower()
    for p in (brand.proof_points if brand and brand.proof_points else []):
        pnum, _ = _metric_parts({"metric": p})
        if pnum and pnum.replace(" ", "").lower() == target:
            return True
        if str(p).strip().lower().startswith(target):
            return True
    return False


def _metric_size(metric: str) -> int:
    """Auto-size the giant number so it always fits the slide width, whatever its length —
    shrink rather than truncate, so a long number (e.g. '$1,250,000') stays correct."""
    n = len(metric)
    if n <= 4:
        return 130
    if n <= 6:
        return 104
    if n <= 9:
        return 80
    if n <= 13:
        return 60
    return 44


def _r_metric(slide, s, idx, total, th):
    _bg(slide, NAVY)
    _rect(slide, 0, 0, Inches(th["spine"]), SH, th["accent"])
    _logo(slide, CREAM)
    _text(slide, Inches(0.9), Inches(1.4), Inches(11.4), Inches(0.8), s.get("title", ""), 28, WHITE, bold=True)
    metric, label = _metric_parts(s)
    _text(slide, Inches(0.85), Inches(2.5), Inches(11.6), Inches(2.2), metric, _metric_size(metric), th["accent"], bold=True)
    _text(slide, Inches(0.95), Inches(5.0), Inches(11), Inches(1.4), label, 26, CREAM)
    _footer(slide, idx, total, CREAM)


def _r_columns(slide, s, idx, total, th, compare=False):
    _bg(slide, WHITE)
    _rect(slide, 0, 0, SW, Inches(1.5), NAVY)
    _rect(slide, 0, 0, Inches(th["spine"]), SH, th["accent"])
    _text(slide, Inches(0.7), Inches(0.55), Inches(11.5), Inches(0.8), s.get("title", ""), 30, WHITE, bold=True)
    left, right = s.get("left") or {}, s.get("right") or {}
    panels = [(Inches(0.7), left, CREAM, NAVY), (Inches(6.95), right, NAVY if compare else CREAM, CREAM if compare else NAVY)]
    for x, col, fill, txt in panels:
        _round(slide, x, Inches(1.95), Inches(5.65), Inches(4.7), fill)
        _text(slide, x + Inches(0.35), Inches(2.2), Inches(5), Inches(0.6), col.get("label", ""), 20, th["accent"] if fill == CREAM else CORAL, bold=True)
        _bullet_list(slide, x + Inches(0.35), Inches(3.0), Inches(5.0), col.get("points", []), txt, size=16, gap=0.7)
    _footer(slide, idx, total)


def _r_quote(slide, s, idx, total, th):
    _bg(slide, CREAM)
    _rect(slide, 0, 0, Inches(th["spine"]), SH, th["accent"])
    _text(slide, Inches(0.8), Inches(0.9), Inches(3), Inches(1.6), "“", 140, th["accent"], bold=True)
    _text(slide, Inches(1.6), Inches(2.4), Inches(11), Inches(3), s.get("quote", s.get("title", "")), 34, NAVY, bold=True)
    if s.get("attribution"):
        _text(slide, Inches(1.65), Inches(5.6), Inches(10), Inches(0.6), f"— {s['attribution']}", 18, NAVY)
    _footer(slide, idx, total)


def _r_closing(slide, s, th):
    _bg(slide, NAVY)
    _rect(slide, 0, 0, SW, Inches(0.16), th["accent"])
    _logo(slide, CREAM)
    _text(slide, Inches(0.9), Inches(2.7), Inches(11.4), Inches(1.6), s.get("title", "Let's talk"), 42, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    _rect(slide, Inches(0.95), Inches(4.5), Inches(2), Inches(0.12), th["accent"])
    _text(slide, Inches(0.9), Inches(4.8), Inches(11), Inches(1), s.get("subtitle") or "Partner with Talentrupt — RPO Done Right.", 22, CREAM)


# --- Public API -----------------------------------------------------------
async def build_deck(
    brand: Brand | None, campaign: Campaign | None, topic: str, slides: int = 6, *,
    audience: str = "", tone: str = "", depth: str = "", design_theme: str = "editorial",
) -> tuple[str, str, dict]:
    slides = max(3, min(slides, 12))
    th = _deck_theme(design_theme)
    outline = await _outline(brand, campaign, topic, slides,
                             audience=audience, tone=tone, depth=depth, design_theme=design_theme)
    cover_img = await _ai_cover(brand, topic)

    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    blank = prs.slide_layouts[6]
    total = len(outline)

    for i, s in enumerate(outline):
        slide = prs.slides.add_slide(blank)
        layout = s.get("layout", "bullets")
        if layout == "cover":
            _r_cover(slide, s, cover_img, th)
        elif layout == "section":
            _r_section(slide, s, th)
        elif layout == "metric" and s.get("metric") and _metric_is_real(s, brand):
            _r_metric(slide, s, i + 1, total, th)
        elif layout == "metric":
            # The number isn't a real proof point — never paint a fabricated stat; show a clean
            # section divider with the intended message instead.
            _r_section(slide, {"title": s.get("title", ""),
                               "subtitle": s.get("metric_label") or s.get("subtitle") or ""}, th)
        elif layout == "two_column":
            _r_columns(slide, s, i + 1, total, th, compare=False)
        elif layout == "comparison":
            _r_columns(slide, s, i + 1, total, th, compare=True)
        elif layout == "cards" and s.get("cards"):
            _r_cards(slide, s, i + 1, total, th)
        elif layout == "timeline":
            _r_timeline(slide, s, i + 1, total, th)
        elif layout == "quote" and (s.get("quote") or s.get("title")):
            _r_quote(slide, s, i + 1, total, th)
        elif layout == "closing":
            _r_closing(slide, s, th)
        else:
            _r_bullets(slide, s, i + 1, total, th)

        # Speaker notes give the presenter talking points — turns a skeletal deck into a usable one.
        note = s.get("notes") or s.get("note")
        if note:
            try:
                slide.notes_slide.notes_text_frame.text = str(note)[:1400]
            except Exception:
                pass

    file_name = unique_name("tr-deck", "pptx")
    path = storage_subdir("decks") / file_name
    prs.save(str(path))
    return str(path), file_name, {
        "topic": topic,
        "slides": total,
        "theme": (design_theme or "editorial"),
        "url": public_url("decks", file_name),
        "ai_cover": bool(cover_img),
    }
