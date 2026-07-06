"""FastAPI entry point — Talentrupt Marketing Agent backend.

Phase 1 surface: health, admin auth, brand, conversations, and SSE chat streaming.
"""
from __future__ import annotations

import asyncio
import csv
import io
import json
import logging
import os
import random
import re
import secrets
import time
import zipfile
from datetime import datetime, timedelta, timezone
from pathlib import Path

from fastapi import BackgroundTasks, Body, Depends, FastAPI, File, Form, Header, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from pydantic import BaseModel
from sqlalchemy import func
from sqlalchemy.orm import Session

from .agent import orchestrator
from .agent.tools import _FEATURE_HEADLINES, _build_one, _save_asset, serialize_asset
from .business import analyze as bd_analyze
from .business import discover as bd_discover
from .business import enrich as bd_enrich
from .business import outreach as bd_outreach
from .business import winstrategy as bd_winstrategy
from .business.profiles import PROFILES
from .business.store import save_opportunity as _save_opp, serialize_opportunity
from .campaigns.planner import interpret_intent, plan_campaign
from .config import settings
from .db import SessionLocal, get_db, init_db
from .generation import decks as gen_decks
from .generation import images as gen_images
from .generation import pdf as gen_pdf
from .generation import posts as gen_posts
from .generation import refine as gen_refine
from .generation import teampost
from .generation.common import public_url, storage_subdir, unique_name
from .knowledge.ingest import ingest_upload, is_supported_upload, run_ingest
from .models import (
    Asset,
    Brand,
    BrandChunk,
    CalendarTask,
    Campaign,
    CampaignItem,
    CampaignProspect,
    Conversation,
    Employee,
    Folder,
    Message,
    Opportunity,
    SourceFile,
)
from .providers import llm
from .schemas import (
    ChatRequest,
    ConversationOut,
    ForgotRequest,
    ForgotResponse,
    LoginRequest,
    LoginResponse,
    MessageOut,
    ResetRequest,
)
from .seed import seed_brand
from . import auth_reset

log = logging.getLogger("talentrupt")

app = FastAPI(title="Talentrupt Marketing Agent", version="2.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=settings.cors_origin_list,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.on_event("startup")
def on_startup() -> None:
    init_db()
    db = SessionLocal()
    try:
        seed_brand(db)
    finally:
        db.close()


# --- Auth -----------------------------------------------------------------
def _role_for_token(token: str) -> str | None:
    """Map a bearer token to its role. Admin and member each have a distinct token."""
    if token and secrets.compare_digest(token, settings.admin_token):
        return "admin"
    if token and secrets.compare_digest(token, settings.member_token):
        return "member"
    return None


def require_auth(authorization: str = Header(default="")) -> str:
    """Validate the bearer token and RETURN the caller's role ('admin' | 'member')."""
    token = authorization.replace("Bearer ", "").strip()
    role = _role_for_token(token)
    if role is None:
        raise HTTPException(status_code=401, detail="Unauthorized")
    return role


def require_admin(role: str = Depends(require_auth)) -> str:
    """Gate admin-only endpoints (Tasks, Analytics). Members get a 403."""
    if role != "admin":
        raise HTTPException(status_code=403, detail="Admins only")
    return role


# --- Health ---------------------------------------------------------------
def _app_version() -> str:
    """The deployed commit SHA, written by CI into _version.txt at package time.

    Lets us confirm from the outside (`curl /api/health`) exactly which commit is live.
    Falls back to "dev" for local runs where the file isn't present.
    """
    try:
        return (Path(__file__).with_name("_version.txt").read_text(encoding="utf-8").strip()
                or "dev")
    except OSError:
        return "dev"


@app.get("/api/health")
def health() -> dict:
    return {
        "status": "ok",
        "version": _app_version(),
        "llm_provider": settings.llm_provider,
        "llm_ready": llm.provider_available(),
        "image_model": settings.openai_image_model,            # configured image model
        "image_model_last": getattr(llm, "LAST_IMAGE_MODEL", settings.openai_image_model),  # model that actually ran last
        "enrichment_ready": settings.enrichment_available(),
        "cutout_ready": _cutout_ready(),        # True once a BG_REMOVAL_API_KEY is live (clean cut-outs on any photo)
        "faceswap_ready": settings.faceswap_available(),
        "database": settings.database_url.split(":")[0],
    }


def _cutout_ready() -> bool:
    try:
        from .generation import cutout
        return cutout.cutout_available()
    except Exception:
        return False


# --- Auth endpoints -------------------------------------------------------
def _is_admin_email(email: str) -> bool:
    return (email or "").strip().lower() == settings.admin_username.strip().lower()


def _is_member_login(username: str, password: str) -> bool:
    return (
        (username or "").strip().lower() == settings.member_username.strip().lower()
        and bool(password)
        and secrets.compare_digest(password, settings.member_password)
    )


@app.post("/api/auth/login", response_model=LoginResponse)
def login(req: LoginRequest, db: Session = Depends(get_db)) -> LoginResponse:
    # Password may have been changed via 'forgot password' (DB override); verify_password falls back
    # to the configured default until then, so existing credentials keep working.
    if _is_admin_email(req.username) and auth_reset.verify_password(db, req.password):
        return LoginResponse(token=settings.admin_token, username=settings.admin_username, role="admin")
    if _is_member_login(req.username, req.password):
        return LoginResponse(token=settings.member_token, username=settings.member_username, role="member")
    raise HTTPException(status_code=401, detail="Invalid credentials")


@app.get("/api/auth/me")
def whoami(role: str = Depends(require_auth)) -> dict:
    """The current session's identity — role is derived from the token (can't be spoofed client-side)."""
    username = settings.admin_username if role == "admin" else settings.member_username
    return {"username": username, "role": role}


@app.post("/api/auth/forgot", response_model=ForgotResponse)
def forgot_password(req: ForgotRequest, db: Session = Depends(get_db)) -> ForgotResponse:
    """Issue a reset code for the admin account. Always returns the same generic message (no account
    enumeration). The code is emailed when SMTP is configured, otherwise logged server-side."""
    generic = "If that email is registered, a reset code has been sent."
    if not _is_admin_email(req.email):
        return ForgotResponse(message=generic, dev_code=None)
    code = auth_reset.create_reset_code(db)
    if auth_reset.email_available():
        try:
            auth_reset.send_reset_email(settings.admin_username, code)
        except Exception:
            log.warning("password reset email failed to send", exc_info=True)
    else:
        log.warning("PASSWORD RESET CODE for %s = %s (email not configured)", req.email, code)
    dev = code if (settings.auth_reset_dev_return_code and not auth_reset.email_available()) else None
    return ForgotResponse(message=generic, dev_code=dev)


@app.post("/api/auth/reset", response_model=LoginResponse)
def reset_password(req: ResetRequest, db: Session = Depends(get_db)) -> LoginResponse:
    """Verify the reset code and set a new password (DB override), then sign the admin in."""
    if not _is_admin_email(req.email) or not auth_reset.verify_reset_code(db, req.code):
        raise HTTPException(status_code=400, detail="Invalid or expired reset code")
    pw = (req.new_password or "").strip()
    if len(pw) < 6:
        raise HTTPException(status_code=400, detail="Password must be at least 6 characters")
    auth_reset.set_password(db, pw)
    auth_reset.clear_reset_code(db)
    return LoginResponse(token=settings.admin_token, username=settings.admin_username)


# --- Brand ----------------------------------------------------------------
@app.get("/api/brand")
def get_brand(db: Session = Depends(get_db), _: None = Depends(require_auth)) -> dict:
    brand = db.query(Brand).first()
    if not brand:
        raise HTTPException(status_code=404, detail="Brand not seeded")
    return {
        "name": brand.name,
        "tagline": brand.tagline,
        "voice": brand.voice,
        "pillars": brand.pillars,
        "services": brand.services,
        "proof_points": brand.proof_points,
        "brand_kit": brand.brand_kit,
    }


# --- Conversations --------------------------------------------------------
@app.get("/api/conversations", response_model=list[ConversationOut])
def list_conversations(
    kind: str | None = None, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    q = db.query(Conversation).filter(Conversation.owner == role)
    if kind:
        q = q.filter(Conversation.kind == kind)
    return q.order_by(Conversation.id.desc()).all()


@app.delete("/api/conversations/{conversation_id}")
def delete_conversation(
    conversation_id: int, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    c = db.get(Conversation, conversation_id)
    if not c or c.owner != role:
        raise HTTPException(status_code=404, detail="Conversation not found")
    db.delete(c)  # cascades to its messages
    db.commit()
    return {"deleted": conversation_id}


@app.get("/api/conversations/{conversation_id}/messages", response_model=list[MessageOut])
def list_messages(
    conversation_id: int, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    conv = db.get(Conversation, conversation_id)
    if not conv or conv.owner != role:
        raise HTTPException(status_code=404, detail="Conversation not found")
    return (
        db.query(Message)
        .filter(Message.conversation_id == conversation_id)
        .order_by(Message.id)
        .all()
    )


@app.post("/api/conversations/{conversation_id}/truncate")
def truncate_conversation(
    conversation_id: int, body: dict = Body(default={}),
    db: Session = Depends(get_db), role: str = Depends(require_auth),
):
    """Delete the LAST `drop` messages of a conversation. Used by the transcript's 'edit message' action,
    which removes the edited turn + everything after it before the client re-sends the edited prompt — so the
    persisted history matches the on-screen edit (and stays consistent on reload). Owner-checked. Counting
    from the back is index-safe: it doesn't matter whether the client shows any synthetic (unpersisted)
    greeting at the front."""
    conv = db.get(Conversation, conversation_id)
    if not conv or conv.owner != role:
        raise HTTPException(status_code=404, detail="Conversation not found")
    try:
        drop = max(0, int(body.get("drop") or 0))
    except (TypeError, ValueError):
        drop = 0
    dropped = 0
    if drop:
        rows = (
            db.query(Message)
            .filter(Message.conversation_id == conversation_id)
            .order_by(Message.id.desc())
            .limit(drop)
            .all()
        )
        for m in rows:
            db.delete(m)
        dropped = len(rows)
        db.commit()
    return {"dropped": dropped}


# --- Chat / Create (SSE streaming) ----------------------------------------
def _stream(req: ChatRequest, db: Session, mode: str, campaign_id: int | None = None,
            role: str = "admin") -> StreamingResponse:
    """Shared SSE handler. mode='chat' -> assistant; mode='create' -> generation;
    mode='campaign' -> internal-campaign studio (campaign_id attaches every asset to that folder).
    Conversations are tagged with the matching kind + owner so each account lists only its own."""
    if req.conversation_id:
        conv = db.get(Conversation, req.conversation_id)
        if not conv or conv.owner != role:  # can't post into another account's thread
            raise HTTPException(status_code=404, detail="Conversation not found")
    else:
        title = req.message.strip()[:60] or "New conversation"
        conv = Conversation(title=title, kind=mode, campaign_id=campaign_id, owner=role)
        db.add(conv)
        db.commit()
        db.refresh(conv)

    db.add(Message(conversation_id=conv.id, role="user", content=req.message))
    db.commit()

    conv_id = conv.id
    conv_title = conv.title
    user_text = req.message
    attachments = req.attachments or []

    async def event_gen():
        yield _sse("meta", {"conversation_id": conv_id, "title": conv_title})
        stream_db = SessionLocal()
        final_text = ""
        err_text = ""
        collected_assets: list = []
        interrupted = False
        try:
            async for ev in orchestrator.run(
                stream_db, conv_id, user_text, mode=mode, attachments=attachments,
                campaign_id=campaign_id, owner=role,
            ):
                if ev["event"] == "done":
                    final_text = ev["data"]
                    yield _sse("done", {"text": final_text})
                elif ev["event"] == "asset":
                    collected_assets.append(ev["data"])
                    yield _sse("asset", ev["data"])
                elif ev["event"] == "error":
                    err_text = ev["data"] or "Something went wrong."
                    yield _sse("error", {"text": err_text})
                elif ev["event"] == "chips":
                    # Tappable quick-pick options for a conversational ask; data is {"items": [...]}.
                    yield _sse("chips", ev["data"])
                else:
                    yield _sse(ev["event"], {"text": ev["data"]})
        except (GeneratorExit, asyncio.CancelledError):
            interrupted = True  # client disconnected mid-stream
            raise
        except Exception as e:  # unhandled error escaping the run loop
            log.warning("chat stream failed (conv %s): %s", conv_id, e, exc_info=True)
            err_text = "The assistant hit an error and couldn't finish that — please try again."
            try:
                yield _sse("error", {"text": err_text})
            except Exception:
                pass
        finally:
            # Persist a TRUTHFUL terminal assistant turn so history isn't corrupted: the answer,
            # else the error, else an interrupted note — but never a blank bubble or an orphan.
            content = final_text or err_text or (
                "⚠️ The response was interrupted before it finished." if interrupted else ""
            )
            if content or collected_assets:
                try:
                    stream_db.add(Message(
                        conversation_id=conv_id, role="assistant",
                        content=content, assets=collected_assets,
                    ))
                    stream_db.commit()
                except Exception:
                    stream_db.rollback()
            stream_db.close()

    return StreamingResponse(event_gen(), media_type="text/event-stream")


@app.post("/api/chat/stream")
async def chat_stream(
    req: ChatRequest, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    return _stream(req, db, mode="chat", role=role)


@app.post("/api/create/stream")
async def create_stream(
    req: ChatRequest, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    return _stream(req, db, mode="create", role=role)


@app.post("/api/campaigns/{campaign_id}/stream")
async def campaign_stream(
    campaign_id: int, req: ChatRequest, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    """Internal-campaign chat: generate posts/visuals/decks/PDFs straight into this campaign folder."""
    c = db.get(Campaign, campaign_id)
    if not c or c.owner != role:
        raise HTTPException(status_code=404, detail="Campaign not found")
    return _stream(req, db, mode="campaign", campaign_id=campaign_id, role=role)


@app.post("/api/chat/attach")
async def chat_attach(
    file: UploadFile = File(...),
    db: Session = Depends(get_db),
    role: str = Depends(require_auth),
):
    """Accept a user-uploaded file, extract its text/caption, embed it into the brand
    library (folder='Uploads'), and return an excerpt the chat can use immediately."""
    name = file.filename or "attachment"
    if not is_supported_upload(name):
        raise HTTPException(
            status_code=415,
            detail="Unsupported file type. Attach a PDF, an image, or a text file (txt, md, csv, json).",
        )
    # Enforce the size cap while streaming so a huge body can't exhaust memory first.
    limit = 25 * 1024 * 1024
    total = 0
    parts: list[bytes] = []
    while chunk := await file.read(1 << 20):
        total += len(chunk)
        if total > limit:
            raise HTTPException(status_code=413, detail="File too large (max 25 MB).")
        parts.append(chunk)
    data = b"".join(parts)
    if not data:
        raise HTTPException(status_code=400, detail="The file is empty.")
    return await ingest_upload(db, name, data, owner=role)


@app.post("/api/knowledge/upload-brand-file")
async def upload_brand_file(
    file: UploadFile = File(...),
    db: Session = Depends(get_db),
    role: str = Depends(require_auth),
):
    """Add a brand asset (PDF/image/text) to the library under folder='Brand Kit' so it IS used for
    grounding (unlike chat attachments, which go to 'Uploads' and are excluded). Mirrors chat_attach."""
    name = file.filename or "brand-asset"
    if not is_supported_upload(name):
        raise HTTPException(
            status_code=415,
            detail="Unsupported file type. Upload a PDF, an image, or a text file (txt, md, csv, json).",
        )
    limit = 25 * 1024 * 1024
    total = 0
    parts: list[bytes] = []
    while chunk := await file.read(1 << 20):
        total += len(chunk)
        if total > limit:
            raise HTTPException(status_code=413, detail="File too large (max 25 MB).")
        parts.append(chunk)
    data = b"".join(parts)
    if not data:
        raise HTTPException(status_code=400, detail="The file is empty.")
    return await ingest_upload(db, name, data, folder="Brand Kit", owner=role)


# --- Folders (employee photo libraries -> real-person posts) ---------------
class FolderCreate(BaseModel):
    name: str = ""


class FolderGenerate(BaseModel):
    topic: str = ""


def _serialize_employee(e: Employee) -> dict:
    return {
        "id": e.id, "folder_id": e.folder_id, "name": e.name, "role": e.role,
        "file_url": e.file_url,
        "created_at": e.created_at.isoformat() if e.created_at else None,
    }


def _serialize_folder(db: Session, f: Folder) -> dict:
    count = db.query(Employee).filter(Employee.folder_id == f.id).count()
    return {
        "id": f.id, "name": f.name, "employee_count": count,
        "created_at": f.created_at.isoformat() if f.created_at else None,
    }


@app.get("/api/folders")
def list_folders(db: Session = Depends(get_db), role: str = Depends(require_auth)):
    rows = db.query(Folder).filter(Folder.owner == role).order_by(Folder.id.desc()).all()
    return [_serialize_folder(db, f) for f in rows]


@app.post("/api/folders")
def create_folder(req: FolderCreate, db: Session = Depends(get_db), role: str = Depends(require_auth)):
    f = Folder(owner=role, name=(req.name or "").strip()[:200] or "Untitled folder")
    db.add(f)
    db.commit()
    db.refresh(f)
    return _serialize_folder(db, f)


@app.delete("/api/folders/{folder_id}")
def delete_folder(folder_id: int, db: Session = Depends(get_db), role: str = Depends(require_auth)):
    f = db.get(Folder, folder_id)
    if not f or f.owner != role:
        raise HTTPException(status_code=404, detail="Folder not found")
    for e in db.query(Employee).filter(Employee.folder_id == folder_id).all():
        try:
            if e.photo_path:
                os.remove(e.photo_path)
        except OSError:
            pass
    db.delete(f)
    db.commit()
    return {"ok": True}


@app.get("/api/folders/{folder_id}/employees")
def list_employees(folder_id: int, db: Session = Depends(get_db), role: str = Depends(require_auth)):
    f = db.get(Folder, folder_id)
    if not f or f.owner != role:
        raise HTTPException(status_code=404, detail="Folder not found")
    rows = db.query(Employee).filter(Employee.folder_id == folder_id).order_by(Employee.id).all()
    return [_serialize_employee(e) for e in rows]


@app.post("/api/folders/{folder_id}/employees")
async def add_employee(
    folder_id: int,
    name: str = Form(...),
    role_title: str = Form(""),
    file: UploadFile = File(...),
    db: Session = Depends(get_db),
    role: str = Depends(require_auth),
):
    f = db.get(Folder, folder_id)
    if not f or f.owner != role:
        raise HTTPException(status_code=404, detail="Folder not found")
    limit = 25 * 1024 * 1024
    total = 0
    parts: list[bytes] = []
    while chunk := await file.read(1 << 20):
        total += len(chunk)
        if total > limit:
            raise HTTPException(status_code=413, detail="Photo too large (max 25 MB).")
        parts.append(chunk)
    data = b"".join(parts)
    if not data:
        raise HTTPException(status_code=400, detail="The photo is empty.")
    ext = (file.filename or "photo.jpg").rsplit(".", 1)[-1].lower()
    if ext not in ("jpg", "jpeg", "png", "webp", "heic", "heif"):
        ext = "jpg"
    fname = unique_name("emp", ext)
    sub = storage_subdir("employees")
    sub.mkdir(parents=True, exist_ok=True)  # the 'employees' storage dir may not exist yet
    path = sub / fname
    with open(path, "wb") as fh:
        fh.write(data)
    e = Employee(
        owner=role, folder_id=folder_id,
        name=(name or "").strip()[:200] or "Employee",
        role=(role_title or "").strip()[:200],
        photo_path=str(path), file_url=public_url("employees", fname),
    )
    db.add(e)
    db.commit()
    db.refresh(e)
    return _serialize_employee(e)


@app.get("/api/employees")
def list_all_employees(db: Session = Depends(get_db), role: str = Depends(require_auth)):
    """Flat list of ALL this account's employees across folders — powers the @ mention picker in Create/Chat."""
    rows = db.query(Employee).filter(Employee.owner == role).order_by(Employee.name).all()
    return [_serialize_employee(e) for e in rows]


@app.delete("/api/employees/{emp_id}")
def delete_employee(emp_id: int, db: Session = Depends(get_db), role: str = Depends(require_auth)):
    e = db.get(Employee, emp_id)
    if not e or e.owner != role:
        raise HTTPException(status_code=404, detail="Employee not found")
    try:
        if e.photo_path:
            os.remove(e.photo_path)
    except OSError:
        pass
    db.delete(e)
    db.commit()
    return {"ok": True}


@app.post("/api/employees/{emp_id}/generate")
async def generate_employee_post(
    emp_id: int, req: FolderGenerate, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    """Generate ONE branded post featuring this employee's REAL photo (AI scene background; the face is
    kept exactly, never AI-generated). Mirrors the chat 'feature_uploaded_person' flow."""
    e = db.get(Employee, emp_id)
    if not e or e.owner != role:
        raise HTTPException(status_code=404, detail="Employee not found")
    try:
        with open(e.photo_path, "rb") as fh:
            raw = fh.read()
    except OSError:
        raise HTTPException(status_code=400, detail="Couldn't read the employee photo.")
    brand = db.query(Brand).first()
    topic = (req.topic or "").strip()
    head, sub = teampost.split_message(topic) if topic else (random.choice(_FEATURE_HEADLINES), "")
    # Use the employee's REAL photo in a branded template. 'magazine'/'split' show the full real photo
    # (never an AI face, no cut-out lib needed) — the AI-scene path was producing random AI people.
    style = random.choice(["magazine", "split"])
    try:
        path, fname, meta = teampost.build_team_image(
            brand, raw, e.name, e.role, head, sub, random.randint(0, 5), style
        )
    except Exception:
        raise HTTPException(status_code=500, detail="Generation failed — please try again.")
    title = (e.name or "Featured") + (f" — {e.role}" if e.role else "")
    a = _save_asset(
        db, None, "image", title[:380],
        body={"person": e.name, "role": e.role, "headline": head, "subline": sub,
              "kind": "team", "folder_id": e.folder_id},
        file_path=path, file_url=meta["url"], meta={**meta, "employee_id": e.id}, owner=role,
    )
    return serialize_asset(a)


# --- Campaigns ------------------------------------------------------------
def _campaign_detail(c: Campaign) -> dict:
    return {
        "id": c.id,
        "name": c.name,
        "type": getattr(c, "type", "external") or "external",
        "goal": c.goal,
        "audience": c.audience,
        "pillar": c.pillar,
        "channels": c.channels,
        "timeline": c.timeline,
        "kpis": c.kpis,
        "strategy": c.strategy,
        "status": c.status,
        "resolved_sector": _campaign_industry(c),  # the sector actually driving client discovery
        "created_at": c.created_at.isoformat() if c.created_at else None,
    }


@app.get("/api/campaigns")
def list_campaigns(
    status: str | None = None, type: str | None = None,
    db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    # Folder rail shows names only (product rule). The Campaigns planner passes
    # status=planning so old/test campaigns don't clutter the view. `type` splits the rail into
    # internal (promote Talentrupt) vs external (client-targeting) folders.
    q = db.query(Campaign).filter(Campaign.owner == role)
    if status:
        q = q.filter(Campaign.status == status)  # explicit (e.g. 'planning' rail, or 'archived' view)
    else:
        q = q.filter(Campaign.status != "archived")  # default: hide archived from the active rails
    if type in ("internal", "external"):
        q = q.filter(Campaign.type == type)
    rows = q.order_by(Campaign.id.desc()).all()
    return [
        {
            "id": c.id,
            "name": c.name,
            "type": getattr(c, "type", "external") or "external",
            "created_at": c.created_at.isoformat() if c.created_at else None,
            "sector": _campaign_industry(c),  # lets the UI open an existing same-sector folder vs duplicate it
        }
        for c in rows
    ]


@app.post("/api/campaigns")
def create_campaign_endpoint(
    payload: dict, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    """Create an internal campaign SHELL (no strategy/prospects) + its chat thread. The campaign
    chat (POST /api/campaigns/{id}/stream) then drives all content generation into this folder, all
    grounded in the campaign's description/brief (stored in `goal`)."""
    name = (payload.get("name") or "Untitled Campaign").strip()[:280] or "Untitled Campaign"
    ctype = payload.get("type") if payload.get("type") in ("internal", "external") else "internal"
    description = (payload.get("description") or "").strip()
    c = Campaign(name=name, type=ctype, goal=description, status="active", owner=role)
    db.add(c)
    db.commit()
    db.refresh(c)
    conv = Conversation(title=name[:60], kind="campaign", campaign_id=c.id, owner=role)
    db.add(conv)
    db.commit()
    db.refresh(conv)
    return {**_campaign_detail(c), "conversation_id": conv.id, "items": [], "assets": []}


@app.get("/api/campaigns/{campaign_id}/messages")
def campaign_messages(
    campaign_id: int, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    """The campaign chat thread (so reopening the folder restores its conversation)."""
    c = db.get(Campaign, campaign_id)
    if not c or c.owner != role:
        return {"conversation_id": None, "messages": []}
    conv = (
        db.query(Conversation).filter(Conversation.campaign_id == campaign_id)
        .order_by(Conversation.id).first()
    )
    if not conv:
        return {"conversation_id": None, "messages": []}
    msgs = (
        db.query(Message).filter(Message.conversation_id == conv.id).order_by(Message.id).all()
    )
    return {
        "conversation_id": conv.id,
        "messages": [{"role": m.role, "content": m.content, "assets": m.assets or []} for m in msgs],
    }


def _serialize_item(it: CampaignItem, db: Session) -> dict:
    asset = db.get(Asset, it.asset_id) if it.asset_id else None
    return {
        "id": it.id,
        "campaign_id": it.campaign_id,
        "scheduled_date": it.scheduled_date.isoformat() if it.scheduled_date else None,
        "channel": it.channel,
        "format": it.format,
        "topic": it.topic,
        "hook": it.hook,
        "status": it.status,
        "asset": serialize_asset(asset) if asset else None,
    }


@app.get("/api/campaigns/{campaign_id}")
def get_campaign(
    campaign_id: int, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    c = db.get(Campaign, campaign_id)
    if not c or c.owner != role:
        raise HTTPException(status_code=404, detail="Campaign not found")
    assets = (
        db.query(Asset)
        .filter(Asset.campaign_id == campaign_id)
        .order_by(Asset.id.desc())
        .all()
    )
    items = (
        db.query(CampaignItem)
        .filter(CampaignItem.campaign_id == campaign_id)
        .order_by(CampaignItem.scheduled_date, CampaignItem.id)
        .all()
    )
    conv = (
        db.query(Conversation).filter(Conversation.campaign_id == campaign_id)
        .order_by(Conversation.id).first()
    )
    return {
        **_campaign_detail(c),
        "conversation_id": conv.id if conv else None,
        "items": [_serialize_item(i, db) for i in items],
        "assets": [serialize_asset(a) for a in assets],
    }


async def _build_planned_campaign(db: Session, brand, brief: dict, start_date=None, owner: str = "admin") -> dict:
    """Plan + persist a campaign (Campaign + dated CampaignItems). Returns its detail."""
    name = (brief.get("name") or brief.get("goal") or "New Campaign").strip()[:280]
    result = await plan_campaign(brand, brief)
    strat = result.get("strategy", {}) or {}
    # Persist the authoritative target sector so this folder's clients stay coherent. Only a
    # vetted value is ever stored/read — never let a stray sector the strategy LLM may emit win.
    sector = (brief.get("sector") or "").strip()
    if sector in _KNOWN_SECTORS:
        strat["sector"] = sector
    else:
        strat.pop("sector", None)
        if sector:
            log.info("campaign plan: ignoring invalid sector %r", sector)

    c = Campaign(
        name=name, owner=owner, goal=brief.get("goal", ""), audience=brief.get("audience", ""),
        pillar=(strat.get("pillars") or [""])[0] if strat.get("pillars") else "",
        channels=brief.get("channels") or ["LinkedIn"], timeline=brief.get("timeframe", "4 weeks"),
        kpis=strat.get("kpis", []), strategy=strat, status="planning",
    )
    db.add(c)
    db.commit()
    db.refresh(c)

    start = datetime.now(timezone.utc)
    if start_date:
        try:
            start = datetime.fromisoformat(str(start_date)[:10]).replace(tzinfo=timezone.utc)
        except ValueError:
            pass
    for it in result.get("items", []):
        db.add(CampaignItem(
            campaign_id=c.id,
            scheduled_date=start + timedelta(days=int(it.get("day_offset", 0))),
            channel=it.get("channel", "LinkedIn"), format=it.get("format", "post"),
            topic=it.get("topic", ""), hook=it.get("hook", ""),
        ))
    db.commit()
    items = (
        db.query(CampaignItem).filter(CampaignItem.campaign_id == c.id)
        .order_by(CampaignItem.scheduled_date, CampaignItem.id).all()
    )
    return {**_campaign_detail(c), "items": [_serialize_item(i, db) for i in items], "assets": []}


@app.post("/api/campaigns/plan")
async def plan_campaign_endpoint(
    payload: dict, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    brand = db.query(Brand).first()
    brief = {
        "name": payload.get("name") or payload.get("goal") or "New Campaign",
        "goal": payload.get("goal", ""),
        "audience": payload.get("audience", ""),
        "sector": payload.get("sector", ""),
        "channels": payload.get("channels") or ["LinkedIn"],
        "timeframe": payload.get("timeframe", "4 weeks"),
    }
    return await _build_planned_campaign(db, brand, brief, payload.get("start_date"), owner=role)


@app.post("/api/campaigns/plan-chat")
async def plan_campaign_chat(
    payload: dict, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    """Conversational intake: interpret the chat → either ask a follow-up or plan + save."""
    messages = payload.get("messages") or []
    brand = db.query(Brand).first()
    intent = await interpret_intent(brand, messages)
    if intent.get("action") == "plan":
        detail = await _build_planned_campaign(db, brand, intent, owner=role)
        return {
            "done": True,
            "reply": f"Done — I've planned “{detail['name']}” with a brief and a "
                     f"{len(detail.get('items', []))}-item content calendar. Opening it now.",
            "campaign": detail,
        }
    return {"done": False, "reply": intent.get("message") or "What's the goal of this campaign, and who's it for?"}


@app.post("/api/campaign-items/{item_id}/generate")
async def generate_campaign_item(
    item_id: int, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    it = db.get(CampaignItem, item_id)
    if not it:
        raise HTTPException(status_code=404, detail="Item not found")
    brand = db.query(Brand).first()
    campaign = db.get(Campaign, it.campaign_id)
    if not campaign or campaign.owner != role:  # the item's campaign must be the caller's
        raise HTTPException(status_code=404, detail="Item not found")
    concept = it.topic or (it.hook or (campaign.name if campaign else "Talentrupt"))

    asset = None
    if it.format == "image":
        rendered = await gen_images.build_images(brand, campaign, concept, count=1)
        if rendered:
            path, _fn, meta = rendered[0]
            asset = _save_asset(db, it.campaign_id, "image", concept,
                                body={"concept": concept}, file_path=path, file_url=meta["url"], meta=meta,
                                owner=role)
    elif it.format == "deck":
        path, _fn, meta = await gen_decks.build_deck(brand, campaign, concept, slides=6)
        asset = _save_asset(db, it.campaign_id, "deck", concept,
                            body={"topic": concept}, file_path=path, file_url=meta["url"], meta=meta,
                            owner=role)
    elif it.format == "pdf":
        path, _fn, meta = gen_pdf.build_pdf(brand, campaign, kind="one-pager")
        asset = _save_asset(db, it.campaign_id, "pdf", concept,
                            body={"kind": "one-pager"}, file_path=path, file_url=meta["url"], meta=meta,
                            owner=role)
    else:  # post (text)
        items = await gen_posts.generate_posts(brand, campaign, count=1,
                                               platform=it.channel or "LinkedIn", angle=concept)
        body = items[0] if items else {"hook": concept}
        asset = _save_asset(db, it.campaign_id, "post", body.get("hook", concept),
                            body=body, meta={"platform": it.channel or "LinkedIn"}, owner=role)

    if asset:
        it.asset_id = asset.id
        it.status = "generated"
        db.commit()
        db.refresh(it)
    return _serialize_item(it, db)


@app.patch("/api/campaign-items/{item_id}")
def update_campaign_item(
    item_id: int, payload: dict, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    """Reschedule a content-calendar item / set its status. Does NOT generate or touch the asset."""
    it = db.get(CampaignItem, item_id)
    _camp = db.get(Campaign, it.campaign_id) if it else None
    if not it or not _camp or _camp.owner != role:
        raise HTTPException(status_code=404, detail="Item not found")
    if payload.get("scheduled_date"):
        try:
            it.scheduled_date = datetime.fromisoformat(
                str(payload["scheduled_date"])[:10]
            ).replace(tzinfo=timezone.utc)
        except ValueError:
            raise HTTPException(status_code=422, detail="Invalid scheduled_date")
    if "status" in payload:
        st = (payload.get("status") or "").strip()
        if st not in ("planned", "generated"):
            raise HTTPException(status_code=422, detail="Unknown status")
        it.status = st
    db.commit()
    db.refresh(it)
    return _serialize_item(it, db)


# --- Campaign target clients (scored prospects per campaign) ---------------
TARGET_CAMPAIGN_PROSPECTS = 6
FILL_COOLDOWN_SECONDS = 180  # after a genuinely DRY pass, don't re-run a paid web search this soon
PARTIAL_RETRY_SECONDS = 20  # after a PARTIAL fill, allow a quick top-up on the next read
MAX_FILL_ROUNDS = 3  # bounded discover retries per fill (caps latency while topping up to TARGET)

# Serialize fills per campaign (one process) so concurrent GET/Done can't double-insert,
# and remember the last attempt so a dry campaign isn't re-searched on every open.
_fill_locks: dict[int, asyncio.Lock] = {}
_last_fill_attempt: dict[int, float] = {}


def _fill_lock(campaign_id: int) -> asyncio.Lock:
    lock = _fill_locks.get(campaign_id)
    if lock is None:
        lock = asyncio.Lock()
        _fill_locks[campaign_id] = lock
    return lock


def _campaign_query(c: Campaign) -> str:
    """Company-search query for prospect discovery. Keep the (often rich) audience for relevance,
    but STEER toward COMPANIES in the campaign's sector — a persona-only audience like
    'HR professionals…' otherwise makes a poor company search."""
    sector = _campaign_industry(c)
    audience = (c.audience or "").strip()
    if sector and audience:
        return f"{audience} — focus on {sector} companies/employers that are actively hiring at volume"[:400]
    if audience:
        return audience
    if sector:
        return f"{sector} companies and employers that are actively hiring at volume"
    return (c.goal or c.name or "companies hiring at volume").strip()


_KNOWN_SECTORS = (
    "Healthcare", "Staffing & Recruiting", "IT & Software",
    "Finance & Fintech", "Corporate / Non-IT",
)

# When a campaign's sector is (re)set, realign its audience to a clean, company-oriented ICP for
# that sector so the folder's discovered clients match — mirrors the quick-start vertical audiences.
SECTOR_DEFAULT_AUDIENCE = {
    "IT & Software": "IT and software companies hiring software engineers, developers, data and product roles at volume",
    "Healthcare": "Healthcare systems, hospitals and clinical/healthcare staffing agencies hiring nurses, allied-health and clinical roles at volume",
    "Staffing & Recruiting": "Staffing and recruiting agencies overloaded with requisitions that need offshore sourcing and recruiting capacity",
    "Finance & Fintech": "Finance and fintech companies hiring engineering, operations, risk and compliance talent at volume",
    "Corporate / Non-IT": "Mid-market corporate employers hiring finance, operations, administrative, customer-success and sales roles at volume",
}

# Keyword sets used BOTH to infer a campaign's sector and to verify a discovered company's
# segment belongs to that sector (the purity check). "Corporate / Non-IT" is deliberately broad
# (a catch-all for non-IT/health/staffing employers) so finance/manufacturing/retail/etc. count.
_SECTOR_KEYWORDS: dict[str, tuple[str, ...]] = {
    "Healthcare": ("healthcare", "health care", "health system", "clinical", "clinic", "nurse",
                   "nursing", "medical", "medicine", "hospital", "patient", "allied health",
                   "physician", "home health", "behavioral health", "pharma", "biotech",
                   "life science", "dental"),
    "Staffing & Recruiting": ("staffing", "recruiting", "recruitment", "search firm",
                              "executive search", "rpo", "talent solutions", "talent acquisition firm"),
    "IT & Software": ("software", "saas", "information technology", "it services", "it consulting",
                      "computer", "internet", "cloud", "developer", "cybersecurity", "technology",
                      "platform", "data infrastructure"),
    "Finance & Fintech": ("fintech", "financial services", "banking", "bank", "insurance",
                          "accounting", "capital markets", "investment", "asset management",
                          "wealth management", "private equity", "venture capital", "credit union",
                          "mortgage", "payments", "lending", "broker"),
    # Corporate / Non-IT is the broad catch-all: a TRUE SUPERSET of Finance plus general non-IT,
    # non-health, non-staffing employers, so finance/manufacturing/retail/etc. all qualify here.
    "Corporate / Non-IT": ("corporate", "manufacturing", "retail", "logistics", "supply chain",
                           "consumer goods", "financial services", "banking", "bank", "insurance",
                           "accounting", "capital markets", "investment", "asset management",
                           "real estate", "hospitality", "education", "construction", "energy",
                           "utilities", "professional services", "business consulting",
                           "telecommunications", "automotive", "transportation", "back office",
                           "restaurant", "aerospace", "defense", "legal", "law firm", "nonprofit",
                           "government", "public sector"),
}

# Word-boundary matching so short tokens don't false-match (e.g. "bank" in "Riverbank",
# "patient" in "Impatient", "hospital" in "Hospitality", "technology" in "Biotechnology").
_SECTOR_RX: dict[str, "re.Pattern[str]"] = {
    s: re.compile(r"\b(?:" + "|".join(re.escape(k) for k in kws) + r")\b")
    for s, kws in _SECTOR_KEYWORDS.items()
}


def _sectors_for_segment(segment: str) -> set[str]:
    """All known sectors whose keywords appear (word-boundary) in this free-text segment label."""
    seg = (segment or "").lower()
    return {s for s, rx in _SECTOR_RX.items() if rx.search(seg)}


def _segment_ok_for_sector(item: dict, campaign_sector: str) -> bool:
    """Purity gate. Keep a company if EITHER the LLM's explicit per-company `sector` classification
    matches the campaign, OR its free-text `segment` clearly belongs to the campaign's sector. The
    second clause matters for OVERLAPPING segments: a "healthcare staffing agency" is a valid target
    for a Healthcare campaign even when the LLM labels its sector "Staffing & Recruiting" — without it
    the gate dropped every healthcare-staffing client and the folder came back empty. For items with
    no classification AND no recognizable segment (ambiguous), keep rather than over-drop."""
    if not campaign_sector:
        return True
    matched = _sectors_for_segment(item.get("segment", ""))
    if campaign_sector in matched:
        return True  # segment text belongs to the campaign's sector (rescues overlaps)
    classified = (item.get("sector") or "").strip()
    if classified in _KNOWN_SECTORS:
        return classified == campaign_sector  # authoritative when the segment didn't already vouch
    return not matched  # no class + segment matches no known sector → ambiguous, don't over-drop


def _campaign_industry(c: Campaign) -> str:
    """The campaign's authoritative target SECTOR so its clients come back coherent (all one
    sector) instead of a random mix. Prefers an explicit stored sector, else infers from text.
    Returns "" when no clear sector is implied."""
    explicit = ((c.strategy or {}).get("sector") or "").strip()
    if explicit in _KNOWN_SECTORS:
        return explicit
    text = " ".join([c.name or "", c.audience or "", c.pillar or "", c.goal or ""]).lower()
    checks = [
        ("Healthcare", ("healthcare", "health system", "clinical", "nurse", "nursing",
                         "medical", "hospital", "patient", "allied health")),
        ("Staffing & Recruiting", ("staffing agenc", "recruiting agenc", "recruitment agenc",
                                   "staffing firm", "staffing and recruiting", "staffing & recruiting")),
        ("IT & Software", ("software", "saas", " it ", "information technology", "engineer",
                           "developer", "tech ", "technology compan", "cloud", "ai/ml")),
        ("Finance & Fintech", ("fintech", "financial services", "banking", "insurance", "accounting firm")),
        ("Corporate / Non-IT", ("non-it", "corporate", "back-office", "back office",
                                "administrative", "operations roles")),
    ]
    for industry, kws in checks:
        if any(k in text for k in kws):
            return industry
    return ""


def _serialize_cp(cp: CampaignProspect) -> dict:
    d = cp.data or {}
    return {
        "id": cp.id,
        "campaign_id": cp.campaign_id,
        "company": cp.company,
        "fit_score": cp.fit_score,
        "segment": d.get("segment", ""),
        "sector": d.get("sector", ""),
        "hiring_signal": d.get("hiring_signal", ""),
        "why_now": d.get("why_now", ""),
        "recommended_service": d.get("recommended_service", ""),
        "contacts": bd_discover.sanitize_contacts(cp.company, d.get("contacts"), d.get("website", "")),
        "timing": d.get("timing", {}),
        "source": d.get("source", ""),
        "company_linkedin": bd_discover.company_linkedin_url(cp.company),
        "status": cp.status,
    }


async def _ensure_campaign_prospects(
    db: Session, campaign: Campaign, *, force: bool = False
) -> list[CampaignProspect]:
    """Fill the campaign's ACTIVE prospects up to TARGET (deduped), returning the new rows.

    Serialized per campaign so concurrent calls can't duplicate; cooldown-gated so a dry
    campaign isn't re-searched on every read. `force` (used by the Done action) bypasses
    the cooldown but still never exceeds TARGET.
    """
    async with _fill_lock(campaign.id):
        existing = (
            db.query(CampaignProspect)
            .filter(CampaignProspect.campaign_id == campaign.id)
            .all()
        )
        active_n = sum(1 for cp in existing if cp.status == "active")
        need = TARGET_CAMPAIGN_PROSPECTS - active_n
        if need <= 0:
            return []
        last = _last_fill_attempt.get(campaign.id)
        if not force and last is not None and (time.monotonic() - last) < FILL_COOLDOWN_SECONDS:
            return []  # recently attempted and still short — don't re-search on every open

        seen_names = [cp.company for cp in existing if cp.company]
        seen_lower = {n.strip().lower() for n in seen_names}
        audience = (campaign.audience or "").strip()
        # Scope to ONE coherent sector so a folder's clients aren't a random mix.
        filters: dict = {}
        industry = _campaign_industry(campaign)
        if industry:
            filters["industry"] = industry
        else:
            log.info("campaign %s: no sector resolved — purity gate inactive", campaign.id)
        if audience:
            filters["keywords"] = audience
            # VIBE the campaign's client search: interpret the audience into structured signals
            # (company size / location / buying-signal) to sharpen discovery. The vetted sector stays
            # authoritative — we only ADD fields the ICP infers, never override the sector or keywords.
            try:
                _icp = await bd_discover.vibe_to_icp(audience)
                for _k in ("company_size", "location", "signal"):
                    if _icp.get(_k) and not filters.get(_k):
                        filters[_k] = _icp[_k]
            except Exception:
                pass
        # Discover in up to MAX_FILL_ROUNDS bounded rounds: the purity gate + name-dedupe shrink
        # each batch, and a single web search often names fewer than asked — so over-fetch AND
        # retry (growing `exclude`) until we hit `need` or run dry, rather than leave a thin folder.
        added: list[CampaignProspect] = []
        dropped: list[str] = []
        all_seen = list(seen_names)
        for _round in range(MAX_FILL_ROUNDS):
            short = need - len(added)
            if short <= 0:
                break
            items = await bd_discover.discover(
                None, _campaign_query(campaign), count=min(short + 6, 12),
                filters=filters or None, exclude=all_seen,
            )
            if not items:
                break  # dry round — stop retrying
            progressed = False
            for it in items:
                name = (it.get("company") or "").strip()
                if not name or name.lower() in seen_lower:
                    continue
                # Purity gate: keep this folder to ONE sector — drop cross-sector intruders
                # (e.g. a pure staffing agency surfacing inside a healthcare campaign).
                if not _segment_ok_for_sector(it, industry):
                    dropped.append(f"{name} ({it.get('sector') or it.get('segment', '')})")
                    seen_lower.add(name.lower())  # don't re-evaluate the same reject next round
                    all_seen.append(name)
                    continue
                seen_lower.add(name.lower())
                all_seen.append(name)
                cp = CampaignProspect(
                    campaign_id=campaign.id, company=name[:300],
                    fit_score=float(it.get("fit_score", 0) or 0), data=it, status="active",
                )
                db.add(cp)
                added.append(cp)
                progressed = True
                if len(added) >= need:
                    break
            if not progressed:
                break  # round produced only dupes/rejects — further rounds unlikely to help
        if dropped:
            log.info("campaign %s [%s]: dropped %d off-sector: %s",
                     campaign.id, industry, len(dropped), "; ".join(dropped))
        # The discover await above is slow; if a concurrent re-target changed the campaign's sector
        # meanwhile, don't commit now-stale clients — abort and let the re-target's next read re-fill.
        db.refresh(campaign)
        if _campaign_industry(campaign) != industry:
            db.rollback()
            log.info("campaign %s: sector changed during fill — aborting stale %s batch", campaign.id, industry)
            return []
        db.commit()
        for cp in added:
            db.refresh(cp)

        # Arm the cooldown based on the OUTCOME (not before the work), so a purity-starved or
        # thin pass doesn't lock the folder for 3 minutes. Full cooldown only on a dry pass;
        # a short retry window on a partial fill so the next read keeps topping up to TARGET.
        now = time.monotonic()
        if not added:
            _last_fill_attempt[campaign.id] = now
        elif len(added) < need:
            _last_fill_attempt[campaign.id] = now - FILL_COOLDOWN_SECONDS + PARTIAL_RETRY_SECONDS
        else:
            _last_fill_attempt[campaign.id] = now  # filled to TARGET — next read early-returns anyway
        return added


@app.get("/api/campaigns/{campaign_id}/prospects")
async def list_campaign_prospects(
    campaign_id: int, status: str = "active",
    db: Session = Depends(get_db), role: str = Depends(require_auth),
):
    """Scored target clients for this campaign. status='active' tops the list up to TARGET
    (cooldown-gated); status='done' returns the worked-through history (newest first, no fill)."""
    c = db.get(Campaign, campaign_id)
    if not c or c.owner != role:
        raise HTTPException(status_code=404, detail="Campaign not found")
    if getattr(c, "type", "external") == "internal":
        return []  # internal campaigns promote Talentrupt itself — no client prospecting
    if status == "done":
        done = (
            db.query(CampaignProspect)
            .filter(CampaignProspect.campaign_id == campaign_id, CampaignProspect.status == "done")
            .order_by(CampaignProspect.id.desc())
            .all()
        )
        return [_serialize_cp(cp) for cp in done]
    await _ensure_campaign_prospects(db, c)
    active = (
        db.query(CampaignProspect)
        .filter(CampaignProspect.campaign_id == campaign_id, CampaignProspect.status == "active")
        .all()
    )
    active.sort(key=lambda cp: cp.fit_score, reverse=True)
    return [_serialize_cp(cp) for cp in active[:TARGET_CAMPAIGN_PROSPECTS]]


@app.post("/api/campaign-prospects/{cp_id}/done")
async def campaign_prospect_done(
    cp_id: int, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    """Mark a client handled and pull in one fresh replacement (idempotent)."""
    cp = db.get(CampaignProspect, cp_id)
    _camp = db.get(Campaign, cp.campaign_id) if cp else None
    if not cp or not _camp or _camp.owner != role:
        raise HTTPException(status_code=404, detail="Prospect not found")
    if cp.status == "done":  # idempotent — a repeated Done must not over-fill
        return {"done_id": cp_id, "replacements": []}
    cp.status = "done"
    db.commit()
    campaign = db.get(Campaign, cp.campaign_id)
    added = await _ensure_campaign_prospects(db, campaign, force=True) if campaign else []
    # Return ALL newly-added clients (a single Done can top up >1 slot) so the UI stays in sync.
    return {"done_id": cp_id, "replacements": [_serialize_cp(cp) for cp in added]}


@app.post("/api/campaign-prospects/{cp_id}/revoke")
def campaign_prospect_revoke(
    cp_id: int, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    """Undo a Done — move a client from the history back into the active list. Does not trigger
    a fill (so it never removes a replacement that was pulled in when the client was marked done)."""
    cp = db.get(CampaignProspect, cp_id)
    _camp = db.get(Campaign, cp.campaign_id) if cp else None
    if not cp or not _camp or _camp.owner != role:
        raise HTTPException(status_code=404, detail="Prospect not found")
    if cp.status == "done":
        cp.status = "active"
        db.commit()
        db.refresh(cp)
    return _serialize_cp(cp)


@app.delete("/api/campaign-prospects/{cp_id}")
def delete_campaign_prospect(
    cp_id: int, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    """Remove a single campaign target client (used to clear an entry from the Done history).
    Active clients are topped back up to TARGET on the next read; this never auto-refills."""
    cp = db.get(CampaignProspect, cp_id)
    _camp = db.get(Campaign, cp.campaign_id) if cp else None
    if not cp or not _camp or _camp.owner != role:
        raise HTTPException(status_code=404, detail="Prospect not found")
    db.delete(cp)
    db.commit()
    return {"deleted": cp_id}


@app.post("/api/campaign-prospects/{cp_id}/strategy")
async def campaign_prospect_strategy(
    cp_id: int, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    """A real, grounded 'how to win this client' strategy. Generated once and cached on the row."""
    cp = db.get(CampaignProspect, cp_id)
    _camp = db.get(Campaign, cp.campaign_id) if cp else None
    if not cp or not _camp or _camp.owner != role:
        raise HTTPException(status_code=404, detail="Prospect not found")
    data = dict(cp.data or {})
    if not data.get("strategy"):
        campaign = db.get(Campaign, cp.campaign_id)
        strat = await bd_winstrategy.win_strategy(data, campaign)
        if strat:
            data["strategy"] = strat
            cp.data = data
            db.commit()
    return {"company": cp.company, "strategy": (cp.data or {}).get("strategy") or {}}


@app.patch("/api/campaigns/{campaign_id}")
def update_campaign(
    campaign_id: int, payload: dict, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    c = db.get(Campaign, campaign_id)
    if not c or c.owner != role:
        raise HTTPException(status_code=404, detail="Campaign not found")
    name = (payload.get("name") or "").strip()
    if name:
        c.name = name[:280]
    if "goal" in payload:  # the internal-campaign brief/description that grounds all generation
        c.goal = (payload.get("goal") or "").strip()
    if "status" in payload:  # soft archive/restore — drops from the 'planning' rail, fully recoverable
        new_status = (payload.get("status") or "").strip()
        if new_status not in ("planning", "archived"):
            raise HTTPException(status_code=422, detail="Unknown status")
        c.status = new_status
    if "sector" in payload:
        sector = (payload.get("sector") or "").strip()
        if sector and sector not in _KNOWN_SECTORS:
            raise HTTPException(status_code=422, detail="Unknown sector")
        strat = dict(c.strategy or {})
        if sector:
            strat["sector"] = sector
            # Realign the audience so this folder's clients actually match the chosen sector.
            c.audience = SECTOR_DEFAULT_AUDIENCE.get(sector, c.audience)
        else:
            strat.pop("sector", None)
        c.strategy = strat
        # Re-target: drop current ACTIVE clients + reset the cooldown so the next read re-fills
        # to the new sector. ('done' history is preserved.)
        db.query(CampaignProspect).filter(
            CampaignProspect.campaign_id == campaign_id,
            CampaignProspect.status == "active",
        ).delete(synchronize_session=False)
        _last_fill_attempt.pop(campaign_id, None)
        log.info("campaign %s re-targeted to sector %r", campaign_id, sector or "(none)")
    db.commit()
    return _campaign_detail(c)


@app.delete("/api/campaigns/{campaign_id}")
def delete_campaign(
    campaign_id: int, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    c = db.get(Campaign, campaign_id)
    if not c or c.owner != role:
        raise HTTPException(status_code=404, detail="Campaign not found")
    # Drop the linked chat thread(s) too (Conversation.campaign_id isn't an ORM cascade).
    for conv in db.query(Conversation).filter(Conversation.campaign_id == campaign_id).all():
        db.delete(conv)
    db.delete(c)  # cascades to assets, prospects, and items (ORM delete-orphan)
    db.commit()
    _last_fill_attempt.pop(campaign_id, None)  # drop per-campaign fill state so a reused id starts clean
    _fill_locks.pop(campaign_id, None)
    return {"deleted": campaign_id}


@app.get("/api/campaigns/{campaign_id}/export")
def export_campaign(
    campaign_id: int, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    c = db.get(Campaign, campaign_id)
    if not c or c.owner != role:
        raise HTTPException(status_code=404, detail="Campaign not found")
    assets = db.query(Asset).filter(Asset.campaign_id == campaign_id).all()

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        # Brief + text content summary
        lines = [f"# {c.name}", ""]
        if c.strategy:
            for k, v in c.strategy.items():
                lines.append(f"## {k}")
                lines.append(", ".join(str(x) for x in v) if isinstance(v, list) else str(v))
                lines.append("")
        zf.writestr("campaign-brief.md", "\n".join(lines))
        for a in assets:
            if a.file_path and Path(a.file_path).exists():
                zf.write(a.file_path, f"{a.type}s/{Path(a.file_path).name}")
            elif a.type == "post":
                b = a.body or {}
                txt = (
                    f"Platform: {b.get('platform','')}\nHook: {b.get('hook','')}\n\n"
                    f"{b.get('caption','')}\n\nCTA: {b.get('cta','')}\n"
                    f"Hashtags: {' '.join(b.get('hashtags', []) or [])}\n"
                )
                zf.writestr(f"posts/post-{a.id}.txt", txt)
    buf.seek(0)
    safe = "".join(ch for ch in c.name if ch.isalnum() or ch in " -_")[:50].strip() or "campaign"
    return StreamingResponse(
        buf, media_type="application/zip",
        headers={"Content-Disposition": f'attachment; filename="{safe}.zip"'},
    )


@app.get("/api/campaigns/{campaign_id}/prospects/export")
def export_campaign_prospects(
    campaign_id: int, status: str = "active",
    db: Session = Depends(get_db), role: str = Depends(require_auth),
):
    """Download a campaign's target clients as CSV. READ-ONLY — never triggers a fill."""
    c = db.get(Campaign, campaign_id)
    if not c or c.owner != role:
        raise HTTPException(status_code=404, detail="Campaign not found")
    q = db.query(CampaignProspect).filter(CampaignProspect.campaign_id == campaign_id)
    if status in ("active", "done"):
        q = q.filter(CampaignProspect.status == status)
    rows = q.order_by(CampaignProspect.fit_score.desc(), CampaignProspect.id.desc()).all()
    cols = [
        "company", "fit_score", "segment", "sector", "hiring_signal", "why_now",
        "recommended_service", "timing_label", "contact_roles", "company_linkedin", "source", "status",
    ]
    sbuf = io.StringIO()
    w = csv.writer(sbuf)
    w.writerow(cols)
    for cp in rows:
        d = _serialize_cp(cp)
        roles = "; ".join(ct.get("role") or "" for ct in (d.get("contacts") or []) if ct.get("role"))
        w.writerow([
            d["company"], d["fit_score"], d["segment"], d["sector"], d["hiring_signal"],
            d["why_now"], d["recommended_service"], (d.get("timing") or {}).get("label", ""),
            roles, d["company_linkedin"], d["source"], d["status"],
        ])
    out = io.BytesIO(sbuf.getvalue().encode("utf-8-sig"))
    safe = "".join(ch for ch in c.name if ch.isalnum() or ch in " -_")[:50].strip() or "campaign"
    return StreamingResponse(
        out, media_type="text/csv",
        headers={"Content-Disposition": f'attachment; filename="{safe}-clients.csv"'},
    )


# --- Assets ---------------------------------------------------------------
@app.get("/api/assets")
def list_assets(
    type: str | None = None,
    general: bool = False,
    db: Session = Depends(get_db),
    role: str = Depends(require_auth),
):
    """List the caller's assets. `general=true` returns ONLY non-campaign assets (the Chat / Create
    'Your generations' gallery) — campaign-specific images live in that campaign's Generated content tab, so
    a football campaign banner never bleeds into the general Chat area."""
    q = db.query(Asset).filter(Asset.owner == role)
    if general:
        q = q.filter(Asset.campaign_id.is_(None))
    if type:
        q = q.filter(Asset.type == type)
    return [serialize_asset(a) for a in q.order_by(Asset.id.desc()).all()]


@app.delete("/api/assets/{asset_id}")
def delete_asset(
    asset_id: int, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    a = db.get(Asset, asset_id)
    if not a or a.owner != role:
        raise HTTPException(status_code=404, detail="Asset not found")
    # Remove the on-disk file too (avoid orphaned files), then the row.
    if a.file_path:
        try:
            p = Path(a.file_path)
            if p.is_file():
                p.unlink()
        except OSError:
            pass  # best-effort; never block the row delete on a filesystem hiccup
    db.delete(a)
    db.commit()
    return {"deleted": asset_id}


@app.post("/api/assets/{asset_id}/regenerate")
async def regenerate_asset_endpoint(
    asset_id: int, payload: dict, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    """Regenerate / refine an existing asset. Saves a NEW asset (original kept) unless replace=true."""
    _orig = db.get(Asset, asset_id)
    if not _orig or _orig.owner != role:  # can't refine another account's asset
        raise HTTPException(status_code=404, detail="Asset not found or not regeneratable")
    a = await gen_refine.regenerate_asset(
        db, asset_id, (payload.get("instruction") or "").strip(), bool(payload.get("replace"))
    )
    if not a:
        raise HTTPException(status_code=404, detail="Asset not found or not regeneratable")
    return serialize_asset(a)


# --- Brand knowledge (source library) -------------------------------------
@app.get("/api/knowledge/status")
def knowledge_status(db: Session = Depends(get_db), _: None = Depends(require_auth)):
    by_folder = dict(
        db.query(SourceFile.folder, func.count(SourceFile.id))
        .group_by(SourceFile.folder)
        .all()
    )
    return {
        "source_files": db.query(SourceFile).count(),
        "brand_chunks": db.query(BrandChunk).count(),
        "by_folder": by_folder,
        "vision_enabled": settings.knowledge_use_vision,
        "zip_present": Path(settings.knowledge_zip_path).exists(),
    }


@app.post("/api/knowledge/import")
async def knowledge_import(
    background: BackgroundTasks, _: None = Depends(require_auth)
):
    background.add_task(run_ingest)
    return {"status": "started"}


# --- Business Dev (lead targeting) ----------------------------------------
def serialize_task(t: CalendarTask) -> dict:
    return {
        "id": t.id,
        "opportunity_id": t.opportunity_id,
        "title": t.title,
        "kind": t.kind,
        "due_at": t.due_at.isoformat() if t.due_at else None,
        "status": t.status,
        "payload": t.payload or {},
        "created_at": t.created_at.isoformat() if t.created_at else None,
    }




@app.get("/api/business/profiles")
def business_profiles(_: None = Depends(require_auth)):
    return PROFILES


@app.post("/api/business/discover")
async def business_discover(
    payload: dict, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    # Exclude companies already in THIS account's pipeline so a repeat search surfaces DIFFERENT firms
    # (and so one account's pipeline never suppresses the other's discoveries).
    known = [c for (c,) in db.query(Opportunity.company).filter(Opportunity.owner == role).all() if c]
    items = await bd_discover.discover(
        payload.get("profile_key"), payload.get("query", ""),
        payload.get("count", 8), payload.get("filters"), exclude=known,
    )
    saved = [serialize_opportunity(_save_opp(db, d, owner=role)) for d in items]
    return {"count": len(saved), "opportunities": saved}


@app.post("/api/business/vibe-discover")
async def business_vibe_discover(
    payload: dict, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    """VIBE PROSPECTING: interpret a freeform 'ideal client' description (the vibe) into a sharp ICP,
    then discover REAL matching companies (fit-scored), ranked by fit. Returns the interpreted ICP +
    the saved prospects so the UI can show 'here's how I read your vibe' + the list."""
    vibe = (payload.get("vibe") or "").strip()
    if not vibe:
        raise HTTPException(status_code=400, detail="Describe your ideal client (the vibe) first.")
    icp = await bd_discover.vibe_to_icp(vibe)
    filters = {k: icp[k] for k in ("industry", "company_size", "location", "signal", "keywords") if icp.get(k)}
    query = icp.get("refined_query") or vibe
    try:
        n = max(1, min(int(payload.get("count", 8) or 8), 12))
    except (TypeError, ValueError):
        n = 8
    # Exclude this account's existing pipeline so a repeat vibe surfaces DIFFERENT firms.
    known = [c for (c,) in db.query(Opportunity.company).filter(Opportunity.owner == role).all() if c]
    items = await bd_discover.discover(None, query, count=n, filters=filters or None, exclude=known)
    saved = [_save_opp(db, d, owner=role) for d in items]
    saved.sort(key=lambda o: -(o.fit_score or 0))
    return {"icp": icp, "count": len(saved), "opportunities": [serialize_opportunity(o) for o in saved]}


@app.post("/api/business/intake")
async def business_intake(
    payload: dict, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    company = (payload.get("company") or "").strip()
    if not company:
        raise HTTPException(status_code=400, detail="company is required")
    d = await bd_analyze.analyze_company(company, payload.get("website", ""))
    if not d:
        raise HTTPException(status_code=502, detail="Analysis unavailable")
    return serialize_opportunity(_save_opp(db, d, owner=role))


@app.get("/api/opportunities")
def list_opportunities(
    status: str | None = None, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    q = db.query(Opportunity).filter(Opportunity.owner == role)
    if status:
        q = q.filter(Opportunity.status == status)
    # Newest first (just-generated clients on top), fit score as the tiebreaker.
    rows = q.order_by(
        Opportunity.created_at.desc(), Opportunity.fit_score.desc(), Opportunity.id.desc()
    ).all()
    return [serialize_opportunity(o) for o in rows]


@app.get("/api/opportunities/export")
def export_opportunities(
    status: str | None = None, saved: str | None = None,
    db: Session = Depends(get_db), role: str = Depends(require_auth),
):
    """Download the prospect list as CSV (read-only). Mirrors the list's status/saved filters."""
    q = db.query(Opportunity).filter(Opportunity.owner == role)
    if status:
        q = q.filter(Opportunity.status == status)
    rows = q.order_by(
        Opportunity.created_at.desc(), Opportunity.fit_score.desc(), Opportunity.id.desc()
    ).all()
    data = [serialize_opportunity(o) for o in rows]
    if saved and saved.strip().lower() in ("1", "true", "yes"):
        data = [d for d in data if d.get("saved")]
    cols = [
        "company", "segment", "fit_score", "status", "saved", "country", "hiring_signal",
        "pain_point", "recommended_service", "why_now", "why_fit", "timing_label",
        "company_linkedin", "website", "source", "decision_maker_roles",
        "sent_at", "replied_at", "meeting_at", "created_at",
    ]
    buf = io.StringIO()
    w = csv.writer(buf)
    w.writerow(cols)
    for d in data:
        why = d.get("why") or {}
        log = why.get("outreach_log") or {}
        roles = "; ".join(c.get("role") or "" for c in (why.get("contacts") or []) if c.get("role"))
        w.writerow([
            d.get("company"), d.get("segment"), d.get("fit_score"), d.get("status"),
            d.get("saved"), d.get("country"), d.get("hiring_signal"), d.get("pain_point"),
            d.get("recommended_service"), why.get("why_now"), why.get("why_fit"),
            (why.get("timing") or {}).get("label"), d.get("company_linkedin"),
            why.get("website"), why.get("source"), roles,
            log.get("sent_at"), log.get("replied_at"), log.get("meeting_at"), d.get("created_at"),
        ])
    out = io.BytesIO(buf.getvalue().encode("utf-8-sig"))  # BOM so Excel opens UTF-8 cleanly
    return StreamingResponse(
        out, media_type="text/csv",
        headers={"Content-Disposition": 'attachment; filename="prospects.csv"'},
    )


@app.post("/api/opportunities/bulk")
def bulk_opportunities(
    payload: dict, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    """Apply one action to many prospects in a single commit. Mirrors the single-item endpoints."""
    ids = [int(i) for i in (payload.get("ids") or []) if str(i).strip().lstrip("-").isdigit()]
    action = (payload.get("action") or "").strip()
    # owner filter: can't mutate/delete the other account's prospects by guessing ids.
    rows = (db.query(Opportunity).filter(Opportunity.id.in_(ids), Opportunity.owner == role).all()
            if ids else [])
    ids = [o.id for o in rows]  # restrict cascade deletes (below) to the caller's own rows
    deleted: list[int] = []
    if action == "delete":
        if ids:
            db.query(CalendarTask).filter(CalendarTask.opportunity_id.in_(ids)).delete(
                synchronize_session=False
            )
        for o in rows:
            deleted.append(o.id)
            db.delete(o)
    elif action == "status" and payload.get("status"):
        for o in rows:
            o.status = payload["status"]
    elif action in ("save", "unsave"):
        for o in rows:
            why = dict(o.why or {})
            why["saved"] = (action == "save")
            o.why = why
    else:
        raise HTTPException(status_code=400, detail="Unknown bulk action")
    db.commit()
    updated = [] if action == "delete" else [serialize_opportunity(o) for o in rows]
    return {"updated": updated, "deleted": deleted}


@app.patch("/api/opportunities/{opp_id}")
def update_opportunity(
    opp_id: int, payload: dict, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    o = db.get(Opportunity, opp_id)
    if not o or o.owner != role:
        raise HTTPException(status_code=404, detail="Opportunity not found")
    if payload.get("status"):
        o.status = payload["status"]
    if "saved" in payload:  # ★ shortlist toggle
        why = dict(o.why or {})
        why["saved"] = bool(payload["saved"])
        o.why = why
    db.commit()
    return serialize_opportunity(o)


@app.delete("/api/opportunities")
def clear_opportunities(db: Session = Depends(get_db), role: str = Depends(require_auth)):
    """Clear UNSAVED prospects (and their follow-up tasks). Saved/shortlisted ones are kept."""
    unsaved = [o for o in db.query(Opportunity).filter(Opportunity.owner == role).all()
               if not (o.why or {}).get("saved")]
    ids = [o.id for o in unsaved]
    if ids:
        db.query(CalendarTask).filter(CalendarTask.opportunity_id.in_(ids)).delete(
            synchronize_session=False
        )
        for o in unsaved:
            db.delete(o)
        db.commit()
    return {"deleted": len(ids)}


@app.delete("/api/opportunities/{opp_id}")
def delete_opportunity(
    opp_id: int, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    o = db.get(Opportunity, opp_id)
    if not o or o.owner != role:
        raise HTTPException(status_code=404, detail="Opportunity not found")
    # Clean up the prospect's follow-up tasks too (avoids orphans + id-reuse cross-talk).
    db.query(CalendarTask).filter(CalendarTask.opportunity_id == opp_id).delete(
        synchronize_session=False
    )
    db.delete(o)
    db.commit()
    return {"deleted": opp_id}


@app.post("/api/opportunities/{opp_id}/outreach")
async def opportunity_outreach(
    opp_id: int, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    o = db.get(Opportunity, opp_id)
    if not o or o.owner != role:
        raise HTTPException(status_code=404, detail="Opportunity not found")
    outreach = await bd_outreach.generate_outreach(o)
    why = dict(o.why or {})
    why["outreach"] = outreach
    o.why = why
    if o.status == "new":
        o.status = "contacted"
    db.commit()
    # Schedule follow-up reminders
    for fu in outreach.get("followups", [])[:2]:
        try:
            days = int(fu.get("day", 3))
        except (TypeError, ValueError):
            days = 3
        db.add(CalendarTask(
            opportunity_id=o.id, owner=role,
            title=f"Follow up with {o.company}",
            kind="followup",
            due_at=datetime.now(timezone.utc) + timedelta(days=days),
            payload={"message": fu.get("message", "")},
        ))
    db.commit()
    db.refresh(o)
    return serialize_opportunity(o)


_PIPELINE_ORDER = ["new", "contacted", "replied", "meeting"]


@app.post("/api/opportunities/{opp_id}/track")
def track_opportunity(
    opp_id: int, payload: dict, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    """Record outreach activity (track-only — the app never sends). Stamps sent/replied/meeting
    dates + notes into why['outreach_log'] and only ADVANCES the pipeline status forward."""
    o = db.get(Opportunity, opp_id)
    if not o or o.owner != role:
        raise HTTPException(status_code=404, detail="Opportunity not found")
    why = dict(o.why or {})
    log = dict(why.get("outreach_log") or {})
    history = list(log.get("history") or [])
    now_iso = datetime.now(timezone.utc).isoformat()
    resolve = lambda v: now_iso if v == "now" else v  # noqa: E731
    advance = payload.get("advance_status", True)
    if "channel" in payload:
        log["channel"] = (payload.get("channel") or "").strip()
    if "notes" in payload:
        log["notes"] = (payload.get("notes") or "").strip()
    cur = o.status if o.status in _PIPELINE_ORDER else "new"
    for field, stage in (("sent_at", "contacted"), ("replied_at", "replied"), ("meeting_at", "meeting")):
        if field in payload:
            val = resolve(payload.get(field))
            log[field] = val
            if val:
                history.append({"event": field[:-3], "at": val, "channel": log.get("channel", "")})
                if advance and _PIPELINE_ORDER.index(stage) > _PIPELINE_ORDER.index(cur):
                    o.status = stage
                    cur = stage
    log["history"] = history
    why["outreach_log"] = log
    o.why = why
    db.commit()
    db.refresh(o)
    return serialize_opportunity(o)


@app.post("/api/opportunities/{opp_id}/enrich")
async def enrich_opportunity(
    opp_id: int, db: Session = Depends(get_db), role: str = Depends(require_auth)
):
    """Fetch REAL, provider-verified contacts on demand (no-op if enrichment is unconfigured).
    Results go to why['verified_contacts'] — NEVER into the sanitized why['contacts']."""
    o = db.get(Opportunity, opp_id)
    if not o or o.owner != role:
        raise HTTPException(status_code=404, detail="Opportunity not found")
    if settings.enrichment_available():
        why = dict(o.why or {})
        verified = await bd_enrich.enrich_contacts(
            o.company or "", why.get("contacts") or [], why.get("website", "")
        )
        why["verified_contacts"] = verified
        o.why = why
        db.commit()
        db.refresh(o)
    return serialize_opportunity(o)


@app.get("/api/tasks")
def list_tasks(
    status: str | None = None, db: Session = Depends(get_db), role: str = Depends(require_admin)
):
    q = db.query(CalendarTask).filter(CalendarTask.owner == role)
    if status:  # optional filter; no param = byte-identical to before
        q = q.filter(CalendarTask.status == status)
    rows = q.order_by(CalendarTask.due_at).all()
    return [serialize_task(t) for t in rows]


@app.patch("/api/tasks/{task_id}")
def update_task(
    task_id: int, payload: dict, db: Session = Depends(get_db), role: str = Depends(require_admin)
):
    """Complete / snooze / reschedule a follow-up task. State rides CalendarTask.status + payload."""
    t = db.get(CalendarTask, task_id)
    if not t or t.owner != role:
        raise HTTPException(status_code=404, detail="Task not found")
    pay = dict(t.payload or {})
    new_status = payload.get("status")
    if new_status in ("pending", "done", "snoozed"):
        t.status = new_status
        if new_status == "done":
            pay["completed_at"] = datetime.now(timezone.utc).isoformat()
    snooze = payload.get("snooze_days")
    if snooze:
        try:
            days = int(snooze)
        except (TypeError, ValueError):
            days = 0
        if days > 0:
            t.due_at = datetime.now(timezone.utc) + timedelta(days=days)
            t.status = "pending"
            pay["snoozed_until"] = t.due_at.isoformat()
    if payload.get("due_at"):
        try:
            t.due_at = datetime.fromisoformat(str(payload["due_at"]).replace("Z", "+00:00"))
        except ValueError:
            pass
    if payload.get("note"):
        pay["note"] = str(payload["note"])
    t.payload = pay
    db.commit()
    db.refresh(t)
    return serialize_task(t)


@app.delete("/api/tasks/{task_id}")
def delete_task(
    task_id: int, db: Session = Depends(get_db), role: str = Depends(require_admin)
):
    """Remove a follow-up task entirely."""
    t = db.get(CalendarTask, task_id)
    if not t or t.owner != role:
        raise HTTPException(status_code=404, detail="Task not found")
    db.delete(t)
    db.commit()
    return {"deleted": task_id}


@app.get("/api/analytics/summary")
def analytics_summary(db: Session = Depends(get_db), role: str = Depends(require_admin)):
    """Read-only pipeline rollup for the Analytics dashboard. Pure aggregation — no writes/LLM/web."""
    opps = db.query(Opportunity).filter(Opportunity.owner == role).all()
    by_status = {s: 0 for s in ("new", "contacted", "replied", "meeting")}
    by_sector: dict[str, int] = {}
    saved = sent = replied = 0
    for o in opps:
        by_status[o.status if o.status in by_status else "new"] += 1
        why = o.why or {}
        if why.get("saved"):
            saved += 1
        log = why.get("outreach_log") or {}
        # "Sent" = genuinely contacted (tracked send, or pipeline advanced past 'new'). A merely
        # DRAFTED outreach (why["outreach"]) is not a send, so it no longer counts here.
        if log.get("sent_at") or o.status in ("contacted", "replied", "meeting"):
            sent += 1
        if log.get("replied_at") or o.status in ("replied", "meeting"):
            replied += 1
        secs = _sectors_for_segment(o.segment or "")
        sector = next((s for s in _KNOWN_SECTORS if s in secs), "Other")
        by_sector[sector] = by_sector.get(sector, 0) + 1

    campaigns = db.query(Campaign).filter(Campaign.owner == role).all()
    planning = sum(1 for c in campaigns if c.status == "planning")
    active_clients = (db.query(CampaignProspect)
                      .join(Campaign, CampaignProspect.campaign_id == Campaign.id)
                      .filter(CampaignProspect.status == "active", Campaign.owner == role).count())

    assets_by_type: dict[str, int] = {}
    for (t,) in db.query(Asset.type).filter(Asset.owner == role).all():
        assets_by_type[t] = assets_by_type.get(t, 0) + 1

    now = datetime.now(timezone.utc)
    overdue = due_soon = pending = 0
    for tk in db.query(CalendarTask).filter(CalendarTask.status == "pending", CalendarTask.owner == role).all():
        pending += 1
        if tk.due_at:
            due = tk.due_at if tk.due_at.tzinfo else tk.due_at.replace(tzinfo=timezone.utc)
            if due < now:
                overdue += 1
            elif due <= now + timedelta(days=7):
                due_soon += 1

    return {
        "opportunities": {
            "by_status": by_status,
            "by_sector": [{"sector": k, "count": v} for k, v in sorted(by_sector.items(), key=lambda x: -x[1])],
            "total": len(opps),
            "saved": saved,
        },
        "outreach": {"sent": sent, "replied": replied},
        "campaigns": {"total": len(campaigns), "planning": planning, "active_clients": active_clients},
        "assets": {"by_type": assets_by_type},
        "tasks": {"overdue": overdue, "due_soon": due_soon, "pending": pending},
    }


# --- File serving (public; dev) -------------------------------------------
_ALLOWED_KINDS = {"images", "decks", "pdfs", "employees"}


def _is_deck_chrome(text: str) -> bool:
    """Return True for brand headers, footers, and slide-number boxes that aren't slide content."""
    import re as _re
    t = text.strip().lower()
    if _re.match(r"^\d+\s*/\s*\d+$", t):  # "3/10" slide counter
        return True
    if "talentrupt" in t and ("rpo done right" in t or t == "tr  talentrupt" or t.startswith("tr ")):
        return True
    if t in {"rpo done right"}:
        return True
    return False


@app.get("/api/files/decks/{file_name}/preview")
def preview_deck(file_name: str, db: Session = Depends(get_db), role: str = Depends(require_auth)):
    """Return slide titles + text extracted from a .pptx file so the UI can show an in-app preview."""
    from pptx import Presentation  # already in requirements; import locally to keep startup lean
    base = storage_subdir("decks").resolve()
    target = (base / file_name).resolve()
    if target.parent != base or not target.exists():  # flat dir — reject traversal & subpaths
        raise HTTPException(status_code=404, detail="Not found")
    asset = db.query(Asset).filter(Asset.file_url.like(f"%/decks/{file_name}")).first()
    if asset and asset.owner != role:  # don't preview another account's deck
        raise HTTPException(status_code=404, detail="Not found")
    try:
        prs = Presentation(str(target))
        slides = []
        for i, slide in enumerate(list(prs.slides), 1):
            texts: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if text and not _is_deck_chrome(text):
                    texts.append(text)
            # First remaining text block = slide title; the rest = body
            title = texts[0] if texts else ""
            body = "\n".join(texts[1:]) if len(texts) > 1 else ""
            slides.append({"slide": i, "title": title, "text": body})
        return {"slide_count": len(list(prs.slides)), "slides": slides}
    except Exception as e:
        log.warning("deck preview failed for %s: %s", file_name, e)  # detail stays server-side
        raise HTTPException(status_code=500, detail="Preview unavailable for this file.")


@app.get("/api/files/{kind}/{file_name}")
def serve_file(kind: str, file_name: str):
    if kind not in _ALLOWED_KINDS:
        raise HTTPException(status_code=404, detail="Not found")
    base = storage_subdir(kind).resolve()
    target = (base / file_name).resolve()
    if target.parent != base or not target.exists():  # flat dir — reject traversal & subpaths
        raise HTTPException(status_code=404, detail="Not found")
    # Force inline so browsers open the file in a new tab rather than auto-downloading it.
    return FileResponse(str(target), headers={"Content-Disposition": f"inline; filename=\"{file_name}\""})


def _sse(event: str, data: dict) -> str:
    return f"event: {event}\ndata: {json.dumps(data)}\n\n"


# --- Serve the exported frontend (single-process deploy) ------------------------------------------
# In production the Next.js app is exported to static files (frontend/out) and served by THIS uvicorn
# process, so the WHOLE app runs as one server (no separate Node process). Mounted LAST so every
# /api/* route registered above takes precedence over the catch-all static handler.
_FRONTEND_DIST = os.environ.get("FRONTEND_DIST") or str(
    Path(__file__).resolve().parents[2] / "frontend" / "out"
)


class _SpaStaticFiles(StaticFiles):
    """Static export server with an SPA fallback: a client-routed path with no matching file is served
    index.html (so deep links / refresh load the app). Real /api 404s are left untouched."""

    async def get_response(self, path, scope):
        from starlette.exceptions import HTTPException as _HX

        try:
            return await super().get_response(path, scope)
        except _HX as exc:
            if exc.status_code == 404 and not path.startswith("api"):
                return await super().get_response("index.html", scope)
            raise


if Path(_FRONTEND_DIST).is_dir():
    app.mount("/", _SpaStaticFiles(directory=_FRONTEND_DIST, html=True), name="frontend")
    log.info("serving exported frontend from %s", _FRONTEND_DIST)
else:
    log.warning("frontend dist not found at %s — running API-only (run `npm run build`)", _FRONTEND_DIST)
